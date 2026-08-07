"""公式テンプレートを使うpptx生成スクリプトの共通部品。

実行目的:
- 表紙、章扉、本文、左右比較、表、カードの組み方を1か所へ集約し、
  複数の説明資料を作っても体裁と配色がずれないようにする。

前提:
- `Template_アルパイン設計事務所公式_2026_16x9.pptx`をレイアウト継承元にする。
- python-pptxを利用できるPythonで、リポジトリルートから実行する。

副作用:
- 本モジュール自体はファイルを作らない。呼び出し側が`Presentation.save()`する。

利用側の責務:
- 資料ごとの文言、スライド構成、掲載可否の判断は呼び出し側が持つ。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Template_アルパイン設計事務所公式_2026_16x9.pptx"

# 創屋様向けご説明資料と同じ配色を使い、資料を並べても同じ資料群に見えるようにする。
INK = RGBColor(0x1F, 0x1F, 0x1F)
BLUE = RGBColor(0x00, 0x5B, 0xAC)
# テンプレートが下端の帯やロゴに使う明るい青。最終ページの帯を自前で引くときに使う。
ACCENT = RGBColor(0x00, 0xA8, 0xE6)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
# フロー図で、1つのレーンだけを塗り分けるときの淡い青。
HIGHLIGHT = RGBColor(0xE8, 0xF2, 0xF9)
CARD_LINE = RGBColor(0xC9, 0xD3, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_LATIN = "Segoe UI"
FONT_EA = "游ゴシック Medium"

LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_TWO = 2
LAYOUT_CHAPTER = 4
LAYOUT_BLANK = 6

# 本文として使える領域。タイトルと下端の帯を避けた範囲。
CONTENT_LEFT = 0.5
CONTENT_RIGHT = 12.83
CONTENT_TOP = 1.5
CONTENT_BOTTOM = 7.03


def apply_font(run, *, size: float, bold: bool = False, color: RGBColor = INK) -> None:
    """英数字と日本語の両方へテンプレートのフォントを設定する。

    python-pptxのfont.nameは欧文（latin）しか設定しないため、日本語が既定フォントの
    ままになる。テンプレート既存スライドとの字面差を避けるため、eaとcsも直接指定する。
    """

    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_LATIN
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        element = rpr.find(qn(tag))
        if element is None:
            element = rpr.makeelement(qn(tag), {})
            rpr.append(element)
        element.set("typeface", FONT_EA)


def set_lines(text_frame, lines: list[tuple[str, dict]], *, space_after: float = 2.0) -> None:
    """段落ごとに文言と書式を指定して流し込む。

    テキストをまとめて代入すると書式が1種類になるため、行単位で段落を作る。
    styleは呼び出しごとに消費するので、辞書を使い回さず都度作る。
    """

    text_frame.word_wrap = True
    for index, (text, style) in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.space_after = Pt(style.pop("space_after", space_after))
        if "align" in style:
            paragraph.alignment = style.pop("align")
        paragraph.line_spacing = style.pop("line_spacing", 1.25)
        # 文中の改行は、生の改行文字ではなく改行要素として書き出す。
        # 生の改行文字はPowerPointで改行として扱われる保証がなく、詰まって表示されうる。
        for position, segment in enumerate(text.split("\n")):
            if position:
                paragraph._p.append(paragraph._p.makeelement(qn("a:br"), {}))
            run = paragraph.add_run()
            run.text = segment
            apply_font(run, **style)


def textbox(slide, x, y, w, h, lines, *, anchor=MSO_ANCHOR.TOP):
    """余白なしのテキストボックスを置き、図形や罫線と左端をそろえる。"""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    set_lines(frame, lines)
    return box


def card(slide, x, y, w, h, *, fill=CARD_BG, line=CARD_LINE):
    """本文カードの下地を置く。装飾ストライプは使わず、淡い塗りと細い枠だけで区切る。"""

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def arrow(slide, x, y, *, size: float = 14.0, color: RGBColor = BLUE):
    """処理の並びをつなぐ矢印。図形ではなく記号で置き、高さのぶれを避ける。"""

    return textbox(
        slide,
        x,
        y,
        0.2,
        0.4,
        [("▶", {"size": size, "bold": True, "color": color, "align": PP_ALIGN.CENTER})],
    )


def drop(shape) -> None:
    """未使用プレースホルダーをXMLごと削除し、編集時に空枠が残らないようにする。"""

    shape._element.getparent().remove(shape._element)


def placeholder(slide, idx):
    """レイアウト由来のプレースホルダーをidxで取り出す。無ければNoneを返す。"""

    for item in slide.placeholders:
        if item.placeholder_format.idx == idx:
            return item
    return None


def fill_heading(slide, title: str, lead: str) -> None:
    """タイトルとリード文を、重ならない位置と大きさで配置する。

    レイアウト既定のままだとタイトルの文字がリード文の行へ食い込むため、
    タイトル枠を少し上へ広げ、リード文をタイトルの下端より下から始める。
    """

    # 位置を上書きするときは4辺すべてを指定する。一部だけ設定すると、レイアウトから
    # 継承していた幅が0のまま書き出され、文字が表示されなくなる。
    title_ph = placeholder(slide, 0)
    title_ph.left = Inches(0.62)
    title_ph.top = Inches(0.38)
    title_ph.width = Inches(11.33)
    title_ph.height = Inches(0.62)
    set_lines(title_ph.text_frame, [(title, {"size": 22, "bold": True, "color": INK})])
    lead_ph = placeholder(slide, 17)
    lead_ph.left = Inches(0.62)
    lead_ph.top = Inches(1.06)
    lead_ph.width = Inches(12.2)
    lead_ph.height = Inches(0.24)
    set_lines(lead_ph.text_frame, [(lead, {"size": 10.5, "bold": True, "color": BLUE})])


def content_slide(prs, title: str, lead: str):
    """タイトルとリード文だけを埋め、本文は個別に組む標準スライドを作る。"""

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    fill_heading(slide, title, lead)
    body = placeholder(slide, 18)
    if body is not None:
        drop(body)
    return slide


def chapter_slide(prs, number: str, title: str, heading: str, bullets: list[str]):
    """章扉。左の青パネルへ章番号と章題、右へこの章で分かることを置く。

    左半分はレイアウト側で青一色に塗られているため、章番号と章題は白文字にする。
    """

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CHAPTER])
    set_lines(placeholder(slide, 2).text_frame, [(number, {"size": 12, "bold": True, "color": WHITE})])
    set_lines(placeholder(slide, 0).text_frame, [(title, {"size": 26, "bold": True, "color": WHITE})])
    set_lines(placeholder(slide, 17).text_frame, [(heading, {"size": 12, "bold": True, "color": BLUE})])
    set_lines(
        placeholder(slide, 18).text_frame,
        [(text, {"size": 12, "color": INK, "space_after": 8}) for text in bullets],
    )
    return slide


def two_column_slide(prs, title, lead, left_head, right_head):
    """左右比較スライド。見出しプレースホルダーだけ埋め、中身は個別に組む。"""

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO])
    fill_heading(slide, title, lead)
    # 左右の見出し帯はレイアウト側で濃色に塗られているため、文字は白で置く。
    set_lines(placeholder(slide, 18).text_frame, [(left_head, {"size": 14, "bold": True, "color": WHITE})])
    set_lines(placeholder(slide, 19).text_frame, [(right_head, {"size": 14, "bold": True, "color": WHITE})])
    for idx in (13, 14):
        unused = placeholder(slide, idx)
        if unused is not None:
            drop(unused)
    return slide


def bullet_lines(items: list[str], *, size: float = 12.0, space_after: float = 7.0):
    """中黒付きの箇条書き段落を作る。"""

    return [(f"・{text}", {"size": size, "color": INK, "space_after": space_after}) for text in items]


def table(
    slide,
    x,
    y,
    w,
    headers,
    rows,
    *,
    col_widths,
    header_size=10.5,
    body_size=10.5,
    row_height=0.38,
    bold_first_column=True,
):
    """見出し行を濃紺で塗った表を置く。行数と列幅から高さを決める。"""

    row_count = len(rows) + 1
    shape = slide.shapes.add_table(
        row_count, len(headers), Inches(x), Inches(y), Inches(w), Inches(row_height * row_count)
    )
    grid = shape.table
    for index, width in enumerate(col_widths):
        grid.columns[index].width = Inches(width)
    for column, text in enumerate(headers):
        cell = grid.cell(0, column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.09)
        set_lines(
            cell.text_frame,
            [(text, {"size": header_size, "bold": True, "color": WHITE, "space_after": 0})],
        )
    for row_index, row in enumerate(rows, start=1):
        for column, text in enumerate(row):
            cell = grid.cell(row_index, column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else CARD_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.09)
            bold = bold_first_column and column == 0
            set_lines(
                cell.text_frame,
                [(text, {"size": body_size, "bold": bold, "color": INK, "space_after": 0})],
            )
    return grid


def note(slide, y: float, text: str, *, size: float = 12.0) -> None:
    """スライド下部の補足。断定できない点や前提条件をここに残す。

    投影して読める大きさを下限とし、本文より一段小さい程度に留める。
    """

    textbox(slide, CONTENT_LEFT, y, 12.33, 0.6, [(text, {"size": size, "color": MUTED, "space_after": 0})])


def closing_slide(prs) -> None:
    """終わりを示すページ。

    テンプレートの最終スライドは謝辞と連絡先が固定文言で入っており、スライド側から
    差し替えられない。本文レイアウトはタイトル枠の装飾が空のまま残るため、白紙から組み、
    他ページと同じ下端の帯だけを引いて END を置く。
    """

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    band = card(slide, 0, 7.36, 13.34, 0.14, fill=ACCENT, line=ACCENT)
    band.line.fill.background()
    textbox(
        slide,
        0.5,
        3.35,
        12.33,
        0.9,
        [("END", {"size": 40, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER})],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_page_numbers(prs) -> None:
    """表紙を除く全ページの右下へ「現在ページ / 総ページ」を入れる。

    レイアウトにはスライド番号枠があるが、python-pptxはスライド追加時にこの枠を
    引き継がないため番号が出ない。また番号枠は総ページ数を持てないので、
    既存の進捗報告資料と同じ「n / N」表記のテキストとして自前で置く。
    """

    total = len(prs.slides._sldIdLst)
    for position, slide in enumerate(prs.slides, start=1):
        if position == 1:
            continue
        textbox(
            slide,
            11.3,
            7.0,
            1.53,
            0.28,
            [
                (
                    f"{position} / {total}",
                    {"size": 10, "color": MUTED, "align": PP_ALIGN.RIGHT, "space_after": 0},
                )
            ],
        )


def open_template():
    """テンプレートを開き、体裁見本の作例スライドを外した状態で返す。

    一覧から外すだけではスライド本体がpptx内に残り、新規スライドとファイル名が
    衝突する。関連付けも解除して、パッケージから完全に落とす。
    """

    if not TEMPLATE.is_file():
        raise FileNotFoundError(f"テンプレートがありません: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.get(qn("r:id")))
        slide_id_list.remove(slide_id)
    return prs


def cover_slide(prs, title: str, subtitle: str, byline: str):
    """表紙。タイトル、サブタイトル、日付・作成者行を埋める。"""

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    set_lines(placeholder(slide, 0).text_frame, [(title, {"size": 40, "bold": True, "color": INK})])
    set_lines(placeholder(slide, 1).text_frame, [(subtitle, {"size": 15, "color": MUTED})])
    set_lines(placeholder(slide, 10).text_frame, [(byline, {"size": 11, "color": MUTED})])
    return slide


def require_new_file(output_path: Path) -> Path:
    """既存ファイルを上書きしないことを、書き出し前に確認する。"""

    resolved = output_path.resolve()
    if resolved.exists():
        raise FileExistsError(f"出力先が既に存在します。上書きしません: {resolved}")
    resolved.parent.mkdir(parents=True, exist_ok=True)
    return resolved
