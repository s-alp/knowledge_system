"""スイムレーン形式のフロー図を、pptxの図形として描く共通部品。

実行目的:
- 業務フロー、処理フロー、導入の進め方を、同じ書式・同じ配色で描けるようにする。
- A3横1枚と16:9スライドの両方へ、同じ定義から出し分ける。

設計方針:
- 途中で切れて読めなくなることを避けるため、1つの流れを1つの面に収める。
- 矢印は「同じレーンを横へ」「同じ列を縦へ」「列と列の間で折れる」の3種類だけにし、
  折れ位置を必ず箱と箱の間の空き列に置く。これにより線どうしの交差が起きない。
- 図形はすべてネイティブ図形で描く。画像を貼らないので、受け取った側で文言を直せる。

副作用:
- 渡されたスライドへ図形を追加する。ファイルの読み書きは行わない。
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
import sys

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SCRIPTS_ROOT = Path(__file__).resolve().parent
if str(SCRIPTS_ROOT) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_ROOT))

from alpine_pptx_kit import (  # noqa: E402
    BLUE,
    CARD_BG,
    HIGHLIGHT,
    CARD_LINE,
    INK,
    MUTED,
    WHITE,
    set_lines,
    textbox,
)


@dataclass(frozen=True)
class Step:
    """フロー内の1つの箱。

    key   : 矢印をつなぐときに使う名前
    lane  : 上から数えたレーン番号
    col   : 左から数えた列番号。列は箱1つぶんの幅を持つ
    label : 箱の中の見出し
    note  : 見出しの下に置く補足。無い場合は空文字
    kind  : process（通常）／decision（分岐）／terminator（始点・終点）
    """

    key: str
    lane: int
    col: int
    label: str
    note: str = ""
    kind: str = "process"


@dataclass(frozen=True)
class Link:
    """矢印1本。labelを付けると分岐の条件を線の近くに書ける。"""

    source: str
    target: str
    label: str = ""


@dataclass
class FlowGeometry:
    """描画面の寸法。A3と16:9で数値だけ差し替える。"""

    left: float
    top: float
    width: float
    height: float
    lane_label_width: float
    columns: int
    box_ratio: float = 0.88
    label_size: float = 11.0
    note_size: float = 8.5
    lane_size: float = 11.0
    link_size: float = 8.5
    boxes: dict = field(default_factory=dict)

    def column_width(self) -> float:
        return (self.width - self.lane_label_width) / self.columns

    def lane_height(self, lane_count: int) -> float:
        return self.height / lane_count

    def box_rect(self, step: Step, lane_count: int) -> tuple[float, float, float, float]:
        """1つの箱の左上座標と大きさを返す。"""

        column_width = self.column_width()
        lane_height = self.lane_height(lane_count)
        box_width = column_width * self.box_ratio
        box_height = lane_height * 0.52
        x = self.left + self.lane_label_width + step.col * column_width + (column_width - box_width) / 2
        y = self.top + step.lane * lane_height + (lane_height - box_height) / 2
        return x, y, box_width, box_height


def _add_arrow_head(connector) -> None:
    """線の終端へ矢印を付ける。python-pptxに設定口が無いためXMLへ直接足す。"""

    line = connector.line._get_or_add_ln()
    tail = line.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    line.append(tail)


def _segment(slide, start: tuple[float, float], end: tuple[float, float], *, head: bool):
    """矢印の1区間を直線で引く。折れ線は区間を分けて描く。"""

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0]),
        Inches(start[1]),
        Inches(end[0]),
        Inches(end[1]),
    )
    connector.line.color.rgb = BLUE
    connector.line.width = Pt(1.5)
    if head:
        _add_arrow_head(connector)
    return connector


def _draw_link(slide, geometry: FlowGeometry, lane_count: int, source: Step, target: Step, label: str) -> None:
    """2つの箱を矢印でつなぐ。

    折れ位置は必ず箱と箱の間の空きに取る。箱の上や中を線が横切らないため、
    レーンをまたいでも線どうしが重ならない。
    """

    sx, sy, sw, sh = geometry.box_rect(source, lane_count)
    tx, ty, tw, th = geometry.box_rect(target, lane_count)
    source_middle = sy + sh / 2
    target_middle = ty + th / 2

    if source.lane == target.lane:
        # 同じレーン。右端から次の箱の左端まで、まっすぐ引く。
        _segment(slide, (sx + sw, source_middle), (tx, target_middle), head=True)
        label_x, label_y = (sx + sw + tx) / 2, source_middle
    elif source.col == target.col:
        # 同じ列。下の箱へ、または上の箱へ、まっすぐ落とす。
        if target.lane > source.lane:
            _segment(slide, (sx + sw / 2, sy + sh), (tx + tw / 2, ty), head=True)
            label_x, label_y = sx + sw / 2, (sy + sh + ty) / 2
        else:
            _segment(slide, (sx + sw / 2, sy), (tx + tw / 2, ty + th), head=True)
            label_x, label_y = sx + sw / 2, (ty + th + sy) / 2
    else:
        # 列もレーンも違う。箱と箱の間で1回だけ折る。
        turn = (sx + sw + tx) / 2
        _segment(slide, (sx + sw, source_middle), (turn, source_middle), head=False)
        _segment(slide, (turn, source_middle), (turn, target_middle), head=False)
        _segment(slide, (turn, target_middle), (tx, target_middle), head=True)
        label_x, label_y = turn, (source_middle + target_middle) / 2

    if label:
        textbox(
            slide,
            label_x - 0.55,
            label_y - 0.24,
            1.1,
            0.22,
            [(label, {"size": geometry.link_size, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER, "space_after": 0})],
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _draw_box(slide, geometry: FlowGeometry, lane_count: int, step: Step) -> None:
    """1つの工程を箱として描く。分岐と始点・終点は形を変えて区別できるようにする。"""

    x, y, width, height = geometry.box_rect(step, lane_count)
    shape_type = {
        "process": MSO_SHAPE.ROUNDED_RECTANGLE,
        "decision": MSO_SHAPE.HEXAGON,
        "terminator": MSO_SHAPE.ROUNDED_RECTANGLE,
    }[step.kind]
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG if step.kind == "decision" else WHITE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.5 if step.kind == "terminator" else 1.0)
    shape.shadow.inherit = False
    if step.kind == "process":
        shape.adjustments[0] = 0.12

    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = [
        (
            step.label,
            {
                "size": geometry.label_size,
                "bold": True,
                "color": INK,
                "align": PP_ALIGN.CENTER,
                "space_after": 1 if step.note else 0,
                "line_spacing": 1.1,
            },
        )
    ]
    if step.note:
        lines.append(
            (
                step.note,
                {
                    "size": geometry.note_size,
                    "color": MUTED,
                    "align": PP_ALIGN.CENTER,
                    "space_after": 0,
                    "line_spacing": 1.05,
                },
            )
        )
    set_lines(frame, lines)


def draw_swimlane(
    slide,
    geometry: FlowGeometry,
    lanes: list[str],
    steps: list[Step],
    links: list[Link],
    *,
    highlight_lane: int | None = None,
) -> None:
    """レーンの帯、工程の箱、矢印をこの順で描く。

    帯を先に描くことで、箱と矢印が必ず帯より前面に来る。
    highlight_laneを指定すると、そのレーンだけ淡く塗り分ける。
    自動で動く範囲のように、一目で区別したいレーンがあるときに使う。
    """

    lane_count = len(lanes)
    lane_height = geometry.lane_height(lane_count)
    for index, name in enumerate(lanes):
        y = geometry.top + index * lane_height
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(geometry.left),
            Inches(y),
            Inches(geometry.width),
            Inches(lane_height),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = HIGHLIGHT if index == highlight_lane else WHITE
        band.line.color.rgb = CARD_LINE
        band.line.width = Pt(0.75)
        band.shadow.inherit = False
        band.text_frame.text = ""

        label = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(geometry.left),
            Inches(y),
            Inches(geometry.lane_label_width),
            Inches(lane_height),
        )
        label.fill.solid()
        label.fill.fore_color.rgb = BLUE
        label.line.color.rgb = BLUE
        label.shadow.inherit = False
        label.text_frame.word_wrap = True
        label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_lines(
            label.text_frame,
            [
                (
                    name,
                    {
                        "size": geometry.lane_size,
                        "bold": True,
                        "color": WHITE,
                        "align": PP_ALIGN.CENTER,
                        "space_after": 0,
                        "line_spacing": 1.1,
                    },
                )
            ],
        )

    by_key = {step.key: step for step in steps}
    for step in steps:
        _draw_box(slide, geometry, lane_count, step)
    for link in links:
        _draw_link(slide, geometry, lane_count, by_key[link.source], by_key[link.target], link.label)
