"""エンドユーザー（導入先の決裁者・管理職）向けのタグ・属性抽出ご説明pptxを生成する。

実行目的:
- 創屋様向けご説明資料（技術者向け）と同じ題材を、導入を判断する立場の方が
  読める粒度へ作り直した説明資料を、公式テンプレートの体裁で生成する。

書き分けの方針:
- API名、内部キー名、実装言語、ソース規模などの開発用語は載せない。
- 客先名、実ファイル名、社内パスなどの実データは載せない。例示は架空値に限定する。
- ライセンス費用交渉、知財、外販条件などの社内・対創屋の論点は載せない。
- 計測していない効果（削減時間・金額）は書かない。
- 画面の説明は実画面のスクリーンショットを使わず、架空データの模擬画面を図形で組む。

前提:
- 公式テンプレートをレイアウト継承元にする（`alpine_pptx_kit`）。
- python-pptxを利用できるPythonで、リポジトリルートから実行する。

副作用:
- 指定した新規pptxを1ファイル作成する。既存ファイルは上書きせず処理を中断する。
"""

from __future__ import annotations

from argparse import ArgumentParser
from pathlib import Path
import sys

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches


SCRIPTS_ROOT = Path(__file__).resolve().parent
if str(SCRIPTS_ROOT) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_ROOT))

from alpine_pptx_kit import (  # noqa: E402
    BLUE,
    CARD_BG,
    INK,
    MUTED,
    ROOT,
    WHITE,
    add_page_numbers,
    arrow,
    bullet_lines,
    card,
    chapter_slide,
    closing_slide,
    content_slide,
    cover_slide,
    note,
    open_template,
    require_new_file,
    set_lines,
    table,
    textbox,
    two_column_slide,
)


DEFAULT_OUTPUT = (
    ROOT
    / "output"
    / "pptx"
    / "20260805_ナレッジシステム_タグ属性抽出_エンドユーザー様向けご説明_r6.pptx"
)


def _mock_search_screen(slide) -> None:
    """タグで絞り込んだ検索画面を、図形だけで組んだ模擬画面として置く。

    実画面のスクリーンショットは客先名・実図番・社内パスが写り込むため使わない。
    値はすべて架空のデモデータとし、後から文言を差し替えられるよう図形で構成する。
    """

    frame_x, frame_y, frame_w, frame_h = 0.6, 1.72, 12.13, 4.8
    card(slide, frame_x, frame_y, frame_w, frame_h, fill=WHITE)

    # 画面上端のタイトルバー。実画面の細部は再現せず、何の画面かだけが分かればよい。
    header = card(slide, frame_x, frame_y, frame_w, 0.46, fill=BLUE, line=BLUE)
    header.text_frame.margin_left = Inches(0.22)
    header.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_lines(
        header.text_frame,
        [("図面検索", {"size": 12, "bold": True, "color": WHITE, "space_after": 0})],
    )

    # 左：絞り込み条件。チェック済みを■、未選択を□で表し、図形数を増やさない。
    panel_x, panel_y = frame_x, frame_y + 0.46
    panel_w, panel_h = 3.05, frame_h - 0.46
    card(slide, panel_x, panel_y, panel_w, panel_h, fill=CARD_BG)
    textbox(
        slide,
        panel_x + 0.25,
        panel_y + 0.22,
        2.6,
        0.28,
        [("絞り込み", {"size": 11.5, "bold": True, "color": BLUE})],
    )
    facets = [
        ("客先", ["■ A社", "□ B社", "□ C社"]),
        ("装置", ["■ ガントリー", "□ 治具", "□ コンベア"]),
        ("材質", ["■ SUS304", "■ SS400", "□ A5052"]),
    ]
    # 見出しと選択肢の実寸から次の位置を決める。固定間隔にすると項目数が増えたときに
    # パネルの下端をはみ出す。
    facet_y = panel_y + 0.62
    for label, options in facets:
        textbox(
            slide,
            panel_x + 0.25,
            facet_y,
            2.6,
            0.22,
            [(label, {"size": 10, "bold": True, "color": MUTED})],
        )
        textbox(
            slide,
            panel_x + 0.25,
            facet_y + 0.26,
            2.6,
            0.235 * len(options),
            [
                (option, {"size": 10.5, "color": INK if option.startswith("■") else MUTED, "space_after": 3})
                for option in options
            ],
        )
        facet_y += 0.26 + 0.235 * len(options) + 0.20

    # 右：絞り込み結果。図番・図面名・付いたタグ・数値が1行で読めることを見せる。
    result_x = frame_x + panel_w + 0.28
    result_w = frame_w - panel_w - 0.56
    textbox(
        slide,
        result_x,
        panel_y + 0.22,
        result_w,
        0.28,
        [("検索結果　4件", {"size": 11.5, "bold": True, "color": INK})],
    )
    results = [
        ("SAMPLE-0001", "ガイドレール", "客先:A社　装置:ガントリー　材質:SUS304", "4.8 kg　/　108 点"),
        ("SAMPLE-0014", "走行ベース", "客先:A社　装置:ガントリー　材質:SS400", "12.3 kg　/　64 点"),
        ("SAMPLE-0027", "ケーブルダクト", "客先:A社　装置:ガントリー　材質:SUS304", "1.2 kg　/　9 点"),
        ("SAMPLE-0033", "昇降フレーム", "客先:A社　装置:ガントリー　材質:SS400　熱処理:焼入れ", "23.6 kg　/　87 点"),
    ]
    row_y = panel_y + 0.62
    for number, name, tags, metrics in results:
        card(slide, result_x, row_y, result_w, 0.8, fill=WHITE)
        textbox(
            slide,
            result_x + 0.2,
            row_y + 0.14,
            2.0,
            0.24,
            [(number, {"size": 10.5, "bold": True, "color": INK})],
        )
        textbox(
            slide,
            result_x + 2.3,
            row_y + 0.14,
            3.0,
            0.24,
            [(name, {"size": 10.5, "color": INK})],
        )
        textbox(
            slide,
            result_x + 0.2,
            row_y + 0.46,
            6.2,
            0.24,
            [(tags, {"size": 9.5, "bold": True, "color": BLUE})],
        )
        textbox(
            slide,
            result_x + 6.6,
            row_y + 0.28,
            1.9,
            0.24,
            [(metrics, {"size": 9.5, "color": MUTED, "align": PP_ALIGN.RIGHT})],
        )
        row_y += 0.93


def build_deck(output_path: Path) -> Path:
    """エンドユーザー向け説明資料を新規生成する。"""

    output_path = require_new_file(output_path)
    prs = open_template()

    # --- 1. 表紙 -----------------------------------------------------------
    cover_slide(
        prs,
        "図面の中身を、探せる情報に。",
        "CADタグ・属性の自動抽出　ご説明資料",
        "2026.08.05　/　株式会社アルパイン設計事務所",
    )

    # --- 2. 要旨 -----------------------------------------------------------
    slide = content_slide(
        prs,
        "この資料でお伝えすること",
        "図面を登録するだけで、探すための手がかりが自動で付きます",
    )
    summary = [
        (
            "01",
            "登録するだけで、\n手がかりが付く",
            "図面ファイルを登録すると、客先・装置・材質・熱処理などの見出しが自動で付きます。"
            "担当者が図面を開いて台帳へ書き写す作業は要りません。",
        ),
        (
            "02",
            "付いた情報には、\n必ず根拠が残る",
            "どの値をもとにそう判断したかが一緒に残ります。"
            "判断できないものは空欄のままにし、推測では埋めません。",
        ),
        (
            "03",
            "過去の図面も、\nそのまま対象にできる",
            "今ある図面をそのまま読み取れます。"
            "ナレッジのために図面を作り直したり、書き足したりする作業はありません。",
        ),
    ]
    for index, (number, head, body) in enumerate(summary):
        x = 0.5 + index * 4.17
        card(slide, x, 1.75, 3.94, 4.25)
        textbox(slide, x + 0.32, 2.05, 3.3, 0.45, [(number, {"size": 24, "bold": True, "color": BLUE})])
        textbox(
            slide,
            x + 0.32,
            2.72,
            3.3,
            1.0,
            [(line, {"size": 15, "bold": True, "color": INK, "space_after": 2}) for line in head.split("\n")],
        )
        textbox(slide, x + 0.32, 3.95, 3.3, 1.9, [(body, {"size": 12, "color": INK})])

    # --- 3. CHAPTER 01 -----------------------------------------------------
    chapter_slide(
        prs,
        "CHAPTER　01",
        "現場で起きていること",
        "Current Issues",
        [
            "図面は増え続けるのに、探す手段は人の記憶のまま",
            "人が台帳へ写す方式が、行き詰まる理由",
            "「探せない」ことが、設計時間と品質の両方を削っている",
        ],
    )

    # --- 4. 課題 -----------------------------------------------------------
    slide = content_slide(
        prs,
        "設計現場で、繰り返し起きていること",
        "どれも「図面の中身が情報になっていない」ことから来ています",
    )
    issues = [
        ("過去の図面が探せない", "似た検討をした記憶はあるが、どこにあるか分からない。結局サーバーを辿るか、当時の担当者に聞くことになる。"),
        ("知らなければ、再利用できない", "使える図面があっても、存在を知らなければ使われない。同じ検討を最初からやり直すことになる。"),
        ("ベテランの記憶に依存している", "どこに何があるかが個人の経験に紐づいている。異動・退職のたびに、組織として探せる範囲が狭くなる。"),
        ("中身は開かないと分からない", "材質・質量・部品点数は、図面を開いて確認するしかない。横断で数えることも、集計することもできない。"),
    ]
    for index, (head, body) in enumerate(issues):
        column, row = index % 2, index // 2
        x = 0.5 + column * 6.25
        y = 1.75 + row * 2.3
        card(slide, x, y, 6.08, 2.05)
        textbox(slide, x + 0.32, y + 0.28, 5.45, 0.35, [(head, {"size": 15, "bold": True, "color": BLUE})])
        textbox(slide, x + 0.32, y + 0.85, 5.45, 1.0, [(body, {"size": 12, "color": INK})])
    note(
        slide,
        6.5,
        "図面そのものは資産として残っています。足りないのは、その中身に外から到達するための手がかりです。",
    )

    # --- 5. なぜ整理だけでは足りないか -------------------------------------
    slide = two_column_slide(
        prs,
        "台帳づくりが続かない理由",
        "人が図面を開いて情報を写す方式は、増え続ける図面に追いつけません",
        "人が写す方式",
        "図面から読み取る方式",
    )
    textbox(
        slide,
        0.9,
        2.75,
        5.3,
        4.1,
        bullet_lines(
            [
                "Excelの台帳へ、図番・客先・材質を手入力する",
                "担当者が要点を個人メモに残す",
                "探すときはフォルダを辿り、記憶と勘で当たりをつける",
            ],
            size=12.5,
        )
        + [("", {"size": 10, "space_after": 10})]
        + [("行き詰まる理由", {"size": 13, "bold": True, "color": BLUE, "space_after": 9})]
        + bullet_lines(
            [
                "図面が増えるたびに、人手の作業も増える",
                "入力の粒度と表記が、人によって変わる",
                "過去の図面までは、遡って入力しきれない",
            ],
            size=12,
        ),
    )
    textbox(
        slide,
        7.15,
        2.75,
        5.3,
        4.1,
        bullet_lines(
            [
                "図面ファイルそのものから、中身を機械的に読み取る",
                "読み取った値の表記をそろえ、検索できる形にする",
                "客先・装置・材質などの見出しを自動で付ける",
            ],
            size=12.5,
        )
        + [("", {"size": 10, "space_after": 10})]
        + [("フォルダの整備は、引き続き効果があります", {"size": 13, "bold": True, "color": BLUE, "space_after": 9})]
        + [
            (
                "階層と名称の規則がそろっていると、客先名や案件名をフォルダから読み取れます。"
                "図面の中身と合わせて判断できるため、見出しの精度が上がります。"
                "整備した規則は、この仕組みの上でそのまま活きます。",
                {"size": 12, "color": INK, "space_after": 0},
            )
        ],
    )

    # --- 6. CHAPTER 02 -----------------------------------------------------
    chapter_slide(
        prs,
        "CHAPTER　02",
        "この仕組みでできること",
        "What It Does",
        [
            "図面1枚から、どんな情報が取り出せるか",
            "自動で付いた情報を、どう信用してよいか",
            "実際の設計図面で確認できている範囲",
        ],
    )

    # --- 7. 3ステップ ------------------------------------------------------
    slide = content_slide(
        prs,
        "何をする仕組みか",
        "図面を渡すと、3つの処理を通って「探せる情報」になります",
    )
    steps = [
        ("STEP 1", "中身を読み取る", "図面ファイルを開いて、書かれている値を機械的に取り出します。元のファイルは読むだけで、書き換えません。"),
        ("STEP 2", "表記をそろえる", "書き方のばらつきをそろえ、決まった項目に整えます。読み取れなかった項目は、空欄のまま残します。"),
        ("STEP 3", "見出しを付ける", "登録されている言葉と照らし合わせ、客先・装置・材質などの見出しを付けます。判断の根拠も一緒に残します。"),
    ]
    for index, (label, head, body) in enumerate(steps):
        x = 0.5 + index * 4.17
        card(slide, x, 1.8, 3.94, 3.15)
        textbox(slide, x + 0.32, 2.08, 3.3, 0.3, [(label, {"size": 11.5, "bold": True, "color": BLUE})])
        textbox(slide, x + 0.32, 2.52, 3.3, 0.45, [(head, {"size": 17, "bold": True, "color": INK})])
        textbox(slide, x + 0.32, 3.22, 3.3, 1.6, [(body, {"size": 12, "color": INK})])
        if index < 2:
            textbox(
                slide,
                x + 3.98,
                3.2,
                0.17,
                0.4,
                [("▶", {"size": 14, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER})],
            )
    textbox(
        slide,
        0.5,
        5.4,
        12.33,
        1.1,
        [
            ("この3つを自動で動かせます。設計者が入力する作業はありません。", {"size": 14, "bold": True, "color": BLUE, "space_after": 8}),
            (
                "読み取りは図面を登録したときに動きます。図面を訂正して登録し直せば、その図面だけを読み取り直します。"
                "人が手で直した見出しは、読み取り直しても残ります。",
                {"size": 12, "color": INK, "space_after": 0},
            ),
        ],
    )

    # --- 8. 図面1枚から出るもの --------------------------------------------
    slide = content_slide(
        prs,
        "図面1枚から、このような情報が付与されます",
        "下記は仕組みを説明するための架空の例示です",
    )
    card(slide, 0.5, 1.75, 6.08, 4.15, fill=WHITE)
    textbox(slide, 0.8, 2.05, 5.5, 0.3, [("付与されるタグ（探すための見出し）", {"size": 14, "bold": True, "color": BLUE})])
    tags = [
        "客先：A社",
        "装置：ガントリー",
        "材質：SUS304",
        "材質：SS400",
        "メーカー：空圧機器メーカーB",
        "熱処理：焼入れ",
        "規格：社内規格C",
    ]
    for index, tag in enumerate(tags):
        column, row = index % 2, index // 2
        x = 0.8 + column * 2.78
        y = 2.62 + row * 0.62
        chip = card(slide, x, y, 2.62, 0.44, fill=CARD_BG, line=BLUE)
        frame = chip.text_frame
        frame.margin_left = Inches(0.11)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_lines(frame, [(tag, {"size": 11, "bold": True, "color": BLUE, "space_after": 0})])
    card(slide, 6.75, 1.75, 6.08, 4.15, fill=WHITE)
    textbox(slide, 7.05, 2.05, 5.5, 0.3, [("記録される属性（図面から取り出した値）", {"size": 14, "bold": True, "color": BLUE})])
    table(
        slide,
        7.05,
        2.62,
        5.5,
        ["項目", "値"],
        [
            ["図面番号", "SAMPLE-0001"],
            ["図面名", "ガイドレール"],
            ["質量", "4.8 kg"],
            ["構成部品", "108 点"],
            ["使われている材質", "8 種類"],
        ],
        col_widths=[2.0, 3.5],
        header_size=10.5,
        body_size=10.5,
    )
    textbox(
        slide,
        0.5,
        6.2,
        12.33,
        0.8,
        [
            (
                "「A社向けで、SUS304を使ったガントリー」といった条件で、図面を開かずに絞り込めます。",
                {"size": 13, "bold": True, "color": INK, "space_after": 6},
            ),
            (
                "質量や部品点数は数値として残るため、複数の図面をまたいだ集計にも使えます。",
                {"size": 11.5, "color": INK, "space_after": 0},
            ),
        ],
    )

    # --- 9. 画面イメージ ----------------------------------------------------
    slide = content_slide(
        prs,
        "タグが付くと、探し方はこう変わります",
        "画面イメージです。表示している値はすべて架空のデモデータです",
    )
    _mock_search_screen(slide)
    note(
        slide,
        6.65,
        "客先・装置・材質などを組み合わせて絞り込めます。フォルダをたどらず、図面を1枚ずつ開かずに候補までたどり着けます。",
    )

    # --- 10. 根拠 ----------------------------------------------------------
    slide = two_column_slide(
        prs,
        "自動で付いた情報を、どう信用してよいか",
        "「なぜそう判断したか」が残る設計にしています",
        "根拠が一緒に残ります",
        "分からないものは、空欄のままにします",
    )
    textbox(
        slide,
        0.9,
        2.75,
        5.3,
        4.1,
        [
            ("見出しには、次の3つが必ず付きます。", {"size": 12.5, "color": INK, "space_after": 12}),
        ]
        + bullet_lines(
            [
                "どの値をもとに判断したか",
                "どれくらい確からしいか",
                "なぜその見出しを付けたか",
            ],
            size=12.5,
        )
        + [
            ("", {"size": 10, "space_after": 10}),
            (
                "確認する人は、元の図面を開かずに「この判断は妥当か」を確かめられます。"
                "誤りが見つかった場合は、その場で直せます。",
                {"size": 12, "color": INK, "space_after": 0},
            ),
        ],
    )
    textbox(
        slide,
        7.15,
        2.75,
        5.3,
        4.1,
        [
            ("確実に判断できない項目は、推測で埋めません。", {"size": 12.5, "color": INK, "space_after": 12}),
        ]
        + bullet_lines(
            [
                "読み取れない項目は、空欄のまま残す",
                "「0」と「分からない」を区別する",
                "候補が複数あるときは、1つに決めない",
            ],
            size=12.5,
        )
        + [
            ("", {"size": 10, "space_after": 10}),
            (
                "誤った値が入ると、探した結果も誤り、確認の手間がかえって増えます。"
                "空欄であれば、足りないことがその場で分かります。",
                {"size": 12, "color": INK, "space_after": 0},
            ),
        ],
    )

    # --- 11. 実データ検証 ---------------------------------------------------
    slide = content_slide(
        prs,
        "動作確認の結果",
        "当社で架空の図面を40件用意し、登録から表示まで通して確認しました",
    )
    table(
        slide,
        0.5,
        1.8,
        12.33,
        ["確認した観点", "結果", "評価"],
        [
            ["最後まで処理できたか", "40件すべて（処理できなかった図面なし）", "落ちずに通る"],
            ["部品の構成を取り出せたか", "40件すべて", "全件で取得"],
            ["質量・重量を取り出せたか", "40件中39件", "ほぼ全件で取得"],
            ["材質を取り出せたか", "40件中34件", "8割超で取得"],
            ["客先・装置の見出しが付いたか", "登録済みの言葉に一致した図面のみ", "登録語彙の数で決まる"],
        ],
        col_widths=[3.6, 5.2, 3.53],
    )
    textbox(
        slide,
        0.5,
        4.4,
        12.33,
        1.7,
        [
            ("読み取りの性能と、見出しの数は別の話です", {"size": 13, "bold": True, "color": BLUE, "space_after": 8}),
            (
                "材質・質量・部品構成は図面そのものに書かれているため、ほぼ全件で取り出せます。"
                "一方、客先や装置の見出しは、あらかじめ登録した社名・装置名と一致したときに付きます。"
                "この確認では語彙を数件しか登録していないため、見出しの数は少なくなっています。",
                {"size": 12, "color": INK, "space_after": 8},
            ),
            (
                "つまり見出しの数は、読み取りの精度ではなく登録した語彙の量で決まります。"
                "運用で使う社名・装置名を登録すれば、その分だけ増えます。",
                {"size": 12, "color": INK, "space_after": 0},
            ),
        ],
    )

    # --- 12. CHAPTER 03 ----------------------------------------------------
    chapter_slide(
        prs,
        "CHAPTER　03",
        "導入にあたっての前提",
        "Prerequisites",
        [
            "ICADを使う構成と、使わない構成",
            "効果を出すために、事前に決めておくこと",
            "この仕組みでできないこと",
        ],
    )

    # --- 13. 必要な環境 ----------------------------------------------------
    slide = content_slide(
        prs,
        "必要な環境",
        "ICADを使う構成と、使わない構成のどちらでも運用できます",
    )
    table(
        slide,
        0.5,
        1.8,
        12.33,
        ["構成", "読み取り方", "ICADの要否", "取り出せる情報"],
        [
            [
                "構成A",
                "ICADの図面をそのまま読み取る",
                "登録する時だけ必要",
                "最も多い（材質・質量・部品構成まで）",
            ],
            [
                "構成B",
                "STEP／DXFに出力した図面を読み取る",
                "不要",
                "中程度（材質・質量は残らないことがある）",
            ],
        ],
        col_widths=[1.3, 4.8, 2.6, 3.63],
    )
    for index, (head, body) in enumerate(
        [
            (
                "ICADが要るのは登録の時だけ",
                "検索・閲覧・見出しの表示にICADは要りません。読み取った内容は保存され、以降は保存された情報を見ます。",
            ),
            (
                "既にお使いのICADを使えます",
                "設計者が使っていない時間帯に読み取りを動かす、といった運用ができます。",
            ),
            (
                "ICADなしでも始められます",
                "STEP／DXFに出力したデータを取り込む構成であれば、ICADなしで運用できます。取り出せる情報は少なくなります。",
            ),
        ]
    ):
        x = 0.5 + index * 4.17
        card(slide, x, 3.6, 3.94, 2.4)
        textbox(slide, x + 0.32, 3.88, 3.3, 0.6, [(head, {"size": 13.5, "bold": True, "color": BLUE})])
        textbox(slide, x + 0.32, 4.58, 3.3, 1.3, [(body, {"size": 11.5, "color": INK})])
    note(
        slide,
        6.3,
        "構成Aと構成Bは併用できます。ふだんは構成Bで取り込み、詳しい情報が必要な図面だけ構成Aで読み取ることもできます。",
    )

    # --- 14. 必要な準備 ----------------------------------------------------
    slide = content_slide(
        prs,
        "効果を出すために、決めておくこと",
        "技術ではなく運用の話です。ここが決まっているかで、精度が大きく変わります",
    )
    preparations = [
        ("自社でよく使う言葉を登録する", "自社の客先名・装置名・案件名を登録します。ここに入っている言葉だけが見出しになります。運用しながら足していけます。"),
        ("フォルダ階層と名称をそろえる", "客先名や案件名をフォルダから読み取れると、見出しの精度が上がります。今後の新規案件から規則をそろえるだけでも効果があります。"),
        ("登録するタイミングを決める", "いつ、誰が、どの図面を登録するかを決めます。出図時に登録する運用にすると、抜けが起きにくくなります。"),
        ("確認する場をつくる", "自動で付いた見出しを人が見て、必要なら直す流れを用意します。最初のうちだけでも確認すると、傾向がつかめます。"),
    ]
    for index, (head, body) in enumerate(preparations):
        column, row = index % 2, index // 2
        x = 0.5 + column * 6.25
        y = 1.8 + row * 2.32
        card(slide, x, y, 6.08, 2.08)
        textbox(
            slide,
            x + 0.32,
            y + 0.28,
            5.45,
            0.35,
            [(f"{index + 1}.　{head}", {"size": 15, "bold": True, "color": BLUE})],
        )
        textbox(slide, x + 0.32, y + 0.85, 5.45, 1.05, [(body, {"size": 12, "color": INK})])
    note(
        slide,
        6.5,
        "どれも一度に完璧にする必要はありません。対象を絞って始め、運用しながら言葉と規則を足していきます。",
    )

    # --- 15. できること・できないこと --------------------------------------
    slide = two_column_slide(
        prs,
        "できること／できないこと",
        "判断していただくために、できないこともあわせてお示しします",
        "できること",
        "できないこと・お約束できないこと",
    )
    textbox(
        slide,
        0.9,
        2.75,
        5.3,
        4.1,
        bullet_lines(
            [
                "図面ファイルから中身を読み取る",
                "表記をそろえて、検索できる形にする",
                "客先・装置・材質などの見出しを自動で付ける",
                "判断の根拠と確からしさを残す",
                "質量・部品点数などを数値として蓄える",
                "過去の図面も、まとめて読み取り直す",
            ],
            size=12.5,
        ),
    )
    textbox(
        slide,
        7.15,
        2.75,
        5.3,
        4.1,
        bullet_lines(
            [
                "図面に書かれていないことを読み取る",
                "すべての項目を100%取り出す",
                "書き方が特殊な図枠を、そのまま読み解く",
                "自動で付いた見出しを、人の確認なしで正とする",
            ],
            size=12.5,
        ),
    )

    # --- 16. 進め方 --------------------------------------------------------
    slide = content_slide(
        prs,
        "導入の進め方",
        "小さく試して、効果を確かめて、広げていきます",
    )
    phases = [
        ("STEP 1", "小さく試す", ["対象の案件・装置を1つに絞る", "その範囲の図面を登録する", "どこまで取り出せるかを見る"]),
        ("STEP 2", "言葉と運用を整える", ["自社の客先名・装置名を登録する", "フォルダ階層と名称をそろえる", "登録と確認の担当を決める"]),
        ("STEP 3", "広げる／効果を測る", ["対象範囲を段階的に広げる", "探す時間の変化を実測する", "測った結果で次の投資を判断する"]),
    ]
    for index, (label, head, items) in enumerate(phases):
        x = 0.5 + index * 4.17
        card(slide, x, 1.8, 3.94, 3.55)
        textbox(slide, x + 0.32, 2.08, 3.3, 0.3, [(label, {"size": 11.5, "bold": True, "color": BLUE})])
        textbox(slide, x + 0.32, 2.52, 3.3, 0.45, [(head, {"size": 17, "bold": True, "color": INK})])
        textbox(
            slide,
            x + 0.32,
            3.22,
            3.3,
            2.0,
            [(f"・{item}", {"size": 12, "color": INK, "space_after": 9}) for item in items],
        )
        if index < 2:
            textbox(
                slide,
                x + 3.98,
                3.4,
                0.17,
                0.4,
                [("▶", {"size": 14, "bold": True, "color": BLUE, "align": PP_ALIGN.CENTER})],
            )
    textbox(
        slide,
        0.5,
        5.8,
        12.33,
        1.0,
        [
            ("まずは対象を絞った試験運用からご相談ください。", {"size": 14, "bold": True, "color": BLUE, "space_after": 8}),
            (
                "対象範囲、登録する図面、確認の進め方をご一緒に決めたうえで、実際のデータで何が取り出せるかをお示しします。",
                {"size": 12, "color": INK, "space_after": 0},
            ),
        ],
    )

    # --- 17. 最終スライド ---------------------------------------------------
    closing_slide(prs)

    add_page_numbers(prs)
    prs.save(str(output_path))
    return output_path


def main() -> int:
    parser = ArgumentParser()
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help="新規作成するpptx。既存ファイルは上書きしない。",
    )
    args = parser.parse_args()
    print(build_deck(args.output))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
