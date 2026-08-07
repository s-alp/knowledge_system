"""CADタグ・属性抽出のフロー図を、A3横1枚ずつのpptx・PDF用として生成する。

実行目的:
- 業務フロー、処理フロー、導入の進め方を、途中で切れない1枚もののフロー図にする。
- 同じ定義を機能概要スライドからも使えるようにし、2つの資料で流れが食い違わないようにする。

読み手:
- 導入先のエンドユーザー（決裁者・管理職）。開発側どうしの資料ではない。

書き分けの方針:
- 関数名、内部キー名、実装言語といった実装用語は載せない。
- 開発会社間の分担や納品範囲は載せない。範囲は「本機能」と「ナレッジシステム本体」の関係で書く。
- 記載内容は`docs/tag_extraction_and_assignment_current_spec_2026-07-29.md`の現行仕様に合わせる。

前提:
- python-pptxを利用できるPythonで、リポジトリルートから実行する。
- 社名ロゴは公式テンプレート内の画像を実行時に取り出して使う。画像を別途持たない。

副作用:
- 指定した新規pptxを1ファイル作成する。既存ファイルは上書きせず処理を中断する。
- ロゴ取り出し用に一時ファイルを作り、生成後に削除する。
"""

from __future__ import annotations

from argparse import ArgumentParser
from dataclasses import dataclass
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZipFile
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPTS_ROOT = Path(__file__).resolve().parent
if str(SCRIPTS_ROOT) not in sys.path:
    sys.path.insert(0, str(SCRIPTS_ROOT))

from alpine_flowchart import FlowGeometry, Link, Step, draw_swimlane  # noqa: E402
from alpine_pptx_kit import (  # noqa: E402
    ACCENT,
    BLUE,
    INK,
    MUTED,
    ROOT,
    TEMPLATE,
    require_new_file,
    set_lines,
    textbox,
)


DEFAULT_OUTPUT = (
    ROOT / "output" / "pptx" / "20260805_ナレッジシステム_タグ属性抽出_フロー図_A3_r1.pptx"
)
# 公式テンプレート内の社名ロゴ。テンプレートを差し替えた場合はここも確認する。
LOGO_MEMBER = "ppt/media/image1.jpg"

A3_WIDTH = 16.54
A3_HEIGHT = 11.69


@dataclass(frozen=True)
class Flow:
    """1枚ぶんのフロー図の定義。"""

    key: str
    title: str
    subtitle: str
    lanes: list[str]
    steps: list[Step]
    links: list[Link]
    columns: int
    footnote: str
    highlight_lane: int | None = None


BUSINESS_FLOW = Flow(
    key="business",
    title="業務フロー　図面を登録してから、探して使うまで",
    subtitle="誰が何をするか。色を付けた中段は、人の操作なしに自動で動く部分です",
    lanes=["設計者", "システム\n（自動）", "確認担当"],
    columns=8,
    highlight_lane=1,
    steps=[
        Step("b1", 0, 0, "図面を作成し\n出図する", "ICAD／STEP／DXF", kind="terminator"),
        Step("b2", 0, 1, "システムへ\n登録する", "出図時の登録が\n抜けにくい"),
        Step("b3", 1, 2, "図面の中身を\n読み取る", "元のファイルは\n変更しない"),
        Step("b4", 1, 3, "項目へ整理して\n属性にする", "決められない項目は\n空欄で残す"),
        Step("b5", 1, 4, "辞書と照合して\nタグを付ける", "根拠と確からしさを\n添える"),
        Step("b6", 2, 5, "自動で付いた\n内容を確認する", "根拠を見て\n妥当性を判断する"),
        Step("b7", 2, 6, "補正する／\n語彙を足す", "補正は読み取り\n直しても残る"),
        Step("b8", 0, 7, "タグで絞り込み\n図面を探す", "客先・装置・材質\nなどで絞る", kind="terminator"),
    ],
    links=[
        Link("b1", "b2"),
        Link("b2", "b3"),
        Link("b3", "b4"),
        Link("b4", "b5"),
        Link("b5", "b6"),
        Link("b6", "b7"),
        Link("b7", "b8"),
    ],
    footnote=(
        "辞書へ語彙を足した後に読み取り直すと、登録済みの図面にも反映されます。"
        "確認と補正は、最初のうちだけでも行うと、どの項目が取りにくいかの傾向がつかめます。"
    ),
)

PROCESS_FLOW = Flow(
    key="process",
    title="処理フロー　1図面が、属性とタグになるまで",
    subtitle="形式による分かれ方と、2D図面・3Dモデルの照合まで。ICADが要るのは読み取りの部分だけです",
    lanes=["ナレッジ\nシステム本体", "ICADが動くPC", "一般PC\n（サーバー）"],
    columns=9,
    steps=[
        Step("p1", 0, 0, "図面を\n登録する", "ICAD／STEP\nDXF", kind="terminator"),
        Step("p2", 1, 1, "形式で\n分かれる", kind="decision"),
        Step("p3", 1, 2, "2D図面を\n読み取る", "ICADのとき\n図枠・注記・寸法"),
        Step("p4", 1, 3, "3Dモデルを\n読み取る", "部品構成・材質\n質量"),
        Step("p5", 2, 2, "STEP／DXF\nを読み取る", "文字・階層\nレイヤー"),
        Step("p6", 2, 4, "項目へ整理して\n属性にする", "どれがどの項目か\n判断する"),
        Step("p7", 2, 5, "辞書と照合して\nタグを付ける", "根拠・確からしさ\n理由を添える"),
        Step("p8", 2, 6, "2Dと3Dを\n照合する", "食い違いは競合\nとして記録"),
        Step("p9", 0, 7, "属性・タグ・\n根拠を保存", "読み取った内容も\n残す"),
        Step("p10", 0, 8, "タグで検索\nできる状態に", "以降の検索に\nICADは不要", kind="terminator"),
    ],
    links=[
        Link("p1", "p2"),
        Link("p2", "p3"),
        Link("p2", "p5"),
        Link("p3", "p4"),
        Link("p4", "p6"),
        Link("p5", "p6"),
        Link("p6", "p7"),
        Link("p7", "p8"),
        Link("p8", "p9"),
        Link("p9", "p10"),
    ],
    footnote=(
        "読み取りは1図面につき1回動かします。図面ごとに独立しているため、失敗した図面だけをやり直せます。"
        "図面を訂正して登録し直すと、その図面だけを読み取り直し、人が手で直した内容は保持されます。"
    ),
)

ROLLOUT_FLOW = Flow(
    key="rollout",
    title="導入の進め方　小さく試して、効果を確かめて、広げる",
    subtitle="最初から全体に広げず、対象を絞って始めます。語彙と規則は運用しながら足していきます",
    lanes=["決める・\n準備する", "動かす", "確かめる"],
    columns=9,
    steps=[
        Step("r1", 0, 0, "対象を1つの\n案件に絞る", "1装置ぶんから\n始める", kind="terminator"),
        Step("r2", 0, 1, "対象範囲の\n図面を集める", "過去の図面も\n対象にできる"),
        Step("r3", 1, 2, "図面を登録して\n読み取る", "自動でタグまで\n付く"),
        Step("r4", 2, 3, "取り出せた\n範囲を見る", "空欄・候補・警告\nを確認する"),
        Step("r5", 0, 4, "客先名・装置名\nを辞書へ登録", "見出しの数は\nここで決まる"),
        Step("r6", 0, 5, "フォルダ階層と\n名称を決める", "今後の新規案件\nからでよい"),
        Step("r7", 1, 6, "読み取り直して\n反映する", "登録済みの図面\nにも反映される"),
        Step("r8", 2, 7, "精度の変化を\n確認する", "付いた見出しの\n増え方を見る"),
        Step("r9", 0, 8, "範囲を広げて\n効果を測る", "探す時間の変化を\n実測する", kind="terminator"),
    ],
    links=[
        Link("r1", "r2"),
        Link("r2", "r3"),
        Link("r3", "r4"),
        Link("r4", "r5"),
        Link("r5", "r6"),
        Link("r6", "r7"),
        Link("r7", "r8"),
        Link("r8", "r9"),
    ],
    footnote=(
        "この一巡を回すと、どの項目が取りにくいか、どの語彙が足りないかが具体的に分かります。"
        "効果を数値で示すには、試験運用の前後で探す時間を実測する必要があります。",
    )[0],
)

FLOWS = [BUSINESS_FLOW, PROCESS_FLOW, ROLLOUT_FLOW]


def slide_geometry(flow: Flow) -> FlowGeometry:
    """16:9スライドへ描くときの寸法。A3より小さいので文字を1段落とす。"""

    return FlowGeometry(
        left=0.5,
        top=1.55,
        width=12.33,
        height=4.95,
        lane_label_width=1.15,
        columns=flow.columns,
        label_size=9.0,
        note_size=7.0,
        lane_size=9.0,
        link_size=7.5,
    )


def _a3_geometry(flow: Flow) -> FlowGeometry:
    """A3横1枚へ描くときの寸法。1枚に収めきることを優先する。"""

    return FlowGeometry(
        left=0.6,
        top=1.95,
        width=15.34,
        height=8.3,
        lane_label_width=1.5,
        columns=flow.columns,
        label_size=11.5,
        note_size=9.0,
        lane_size=11.5,
        link_size=9.5,
    )


def _a3_frame(slide, flow: Flow, logo_path: Path) -> None:
    """A3用紙の見出し、ロゴ、下端の帯、注記を置く。"""

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.28), Inches(A3_HEIGHT))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(A3_HEIGHT - 0.16), Inches(A3_WIDTH), Inches(0.16)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    textbox(slide, 0.6, 0.62, 12.5, 0.55, [(flow.title, {"size": 24, "bold": True, "color": INK})])
    textbox(slide, 0.6, 1.32, 12.5, 0.32, [(flow.subtitle, {"size": 11.5, "bold": True, "color": BLUE})])
    slide.shapes.add_picture(str(logo_path), Inches(13.5), Inches(0.55), width=Inches(2.3))

    textbox(
        slide,
        0.6,
        10.5,
        15.34,
        0.6,
        [(flow.footnote, {"size": 11, "color": INK, "space_after": 0})],
    )
    textbox(
        slide,
        0.6,
        11.15,
        15.34,
        0.28,
        [
            (
                "CADタグ・属性抽出　フロー図　/　2026.08.05　株式会社アルパイン設計事務所",
                {"size": 9, "color": MUTED, "space_after": 0},
            )
        ],
    )


def add_flow_slide(prs, flow: Flow, geometry: FlowGeometry) -> None:
    """既に用意されたスライドサイズへ、レーンと箱と矢印を描く。"""

    draw_swimlane(prs.slides[-1], geometry, flow.lanes, flow.steps, flow.links)


def build_a3_deck(output_path: Path) -> Path:
    """A3横1枚ずつのフロー図を、3枚のpptxとして生成する。"""

    output_path = require_new_file(output_path)
    if not TEMPLATE.is_file():
        raise FileNotFoundError(f"テンプレートがありません: {TEMPLATE}")

    prs = Presentation()
    prs.slide_width = Inches(A3_WIDTH)
    prs.slide_height = Inches(A3_HEIGHT)
    blank_layout = prs.slide_layouts[6]

    with TemporaryDirectory() as temp:
        logo_path = Path(temp) / "logo.jpg"
        with ZipFile(TEMPLATE) as archive:
            logo_path.write_bytes(archive.read(LOGO_MEMBER))
        for flow in FLOWS:
            slide = prs.slides.add_slide(blank_layout)
            _a3_frame(slide, flow, logo_path)
            draw_swimlane(
                slide,
                _a3_geometry(flow),
                flow.lanes,
                flow.steps,
                flow.links,
                highlight_lane=flow.highlight_lane,
            )
        prs.save(str(output_path))
    return output_path


def main() -> int:
    parser = ArgumentParser()
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help="新規作成するA3横のpptx。既存ファイルは上書きしない。",
    )
    args = parser.parse_args()
    print(build_a3_deck(args.output))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
