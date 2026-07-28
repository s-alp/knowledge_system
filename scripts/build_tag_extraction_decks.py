"""公式テンプレートから、創屋様向け／役員向けの説明pptxを生成する。

- テンプレート: Template_アルパイン設計事務所公式_2026_16x9.pptx
- 出力先: output/pptx/
- 事実の出典は docs/cad_tag_extraction_sources_for_souya_2026-07-28.md に対応させる。
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Template_アルパイン設計事務所公式_2026_16x9.pptx"
OUTDIR = ROOT / "output" / "pptx"

JP_FONT = "游ゴシック Medium"
EN_FONT = "Segoe UI"

INK = RGBColor(0x1F, 0x1F, 0x1F)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT = RGBColor(0x00, 0x5B, 0xAC)
ACCENT_BG = RGBColor(0xEA, 0xF2, 0xFA)
WARN_BG = RGBColor(0xFD, 0xF3, 0xE3)
LINE = RGBColor(0xC9, 0xD3, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_TWO = 2
LAYOUT_CHAPTER = 4
LAYOUT_END = 7

BODY_LEFT = Emu(457200)
BODY_TOP = Emu(1364228)
BODY_WIDTH = Emu(11352213)
BODY_HEIGHT = Emu(5069910)


# --------------------------------------------------------------------------- utils
def clear_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rid = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rid)
        sld_id_lst.remove(sld_id)


def drop_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def drop_unused_placeholders(slide, used_idx: set[int]) -> None:
    """使わなかったプレースホルダーをXMLレベルで削除する（枠が残らないように）。"""
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx in used_idx:
            continue
        if ph.placeholder_format.type == 13:  # スライド番号は残す
            continue
        drop_shape(ph)


def style_runs(text_frame, *, size=None, bold=None, color=None) -> None:
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.name = JP_FONT
            if size is not None:
                run.font.size = size
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color


def set_bullets(text_frame, items, *, base_size=Pt(16)) -> None:
    """items: [(level, text, bold?)] を流し込む。"""
    text_frame.clear()
    text_frame.word_wrap = True
    first = True
    for item in items:
        level, text = item[0], item[1]
        bold = item[2] if len(item) > 2 else False
        para = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        para.level = level
        run = para.add_run()
        run.text = text
        run.font.name = JP_FONT
        run.font.size = Pt(base_size.pt - level * 2) if level else base_size
        run.font.bold = bold
        run.font.color.rgb = INK if level == 0 else MUTED
        para.space_after = Pt(6 if level == 0 else 3)


# --------------------------------------------------------------------------- slides
def add_title_slide(prs, title, subtitle, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    slide.placeholders[0].text_frame.text = title
    style_runs(slide.placeholders[0].text_frame, size=Pt(36), bold=True, color=INK)
    slide.placeholders[1].text_frame.text = subtitle
    style_runs(slide.placeholders[1].text_frame, size=Pt(18), color=MUTED)
    slide.placeholders[10].text_frame.text = meta
    style_runs(slide.placeholders[10].text_frame, size=Pt(12), color=MUTED)
    return slide


def add_chapter_slide(prs, no, title, lead, intro_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CHAPTER])
    slide.placeholders[2].text_frame.text = f"CHAPTER　{no}"
    style_runs(slide.placeholders[2].text_frame, size=Pt(14), bold=True, color=ACCENT)
    slide.placeholders[0].text_frame.text = title
    style_runs(slide.placeholders[0].text_frame, size=Pt(28), bold=True, color=INK)
    slide.placeholders[17].text_frame.text = lead
    style_runs(slide.placeholders[17].text_frame, size=Pt(12), bold=True, color=ACCENT)
    set_bullets(slide.placeholders[18].text_frame, intro_lines, base_size=Pt(14))
    return slide


def new_content_slide(prs, title, lead):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    slide.placeholders[0].text_frame.text = title
    style_runs(slide.placeholders[0].text_frame, size=Pt(24), bold=True, color=INK)
    slide.placeholders[17].text_frame.text = lead
    style_runs(slide.placeholders[17].text_frame, size=Pt(11), bold=True, color=ACCENT)
    return slide


def add_bullet_slide(prs, title, lead, items, base_size=Pt(16)):
    slide = new_content_slide(prs, title, lead)
    set_bullets(slide.placeholders[18].text_frame, items, base_size=base_size)
    return slide


def add_table_slide(prs, title, lead, header, rows, *, col_widths=None,
                    font_size=Pt(11), header_size=Pt(11), note=None,
                    top=None, height=None, highlight_rows=None):
    slide = new_content_slide(prs, title, lead)
    drop_shape(slide.placeholders[18])

    n_cols = len(header)
    n_rows = len(rows) + 1
    tbl_top = top or BODY_TOP
    tbl_height = height or Emu(int(BODY_HEIGHT * 0.9))
    graphic = slide.shapes.add_table(
        n_rows, n_cols, BODY_LEFT, tbl_top, BODY_WIDTH, tbl_height
    )
    table = graphic.table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(int(BODY_WIDTH * w / total))

    highlight_rows = highlight_rows or set()

    for c, text in enumerate(header):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Emu(64000)
        cell.margin_right = Emu(64000)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = JP_FONT
                run.font.size = header_size
                run.font.bold = True
                run.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(text)
            cell.fill.solid()
            if r in highlight_rows:
                cell.fill.fore_color.rgb = ACCENT_BG
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF5, 0xF7, 0xFA)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Emu(64000)
            cell.margin_right = Emu(64000)
            cell.margin_top = Emu(24000)
            cell.margin_bottom = Emu(24000)
            for para in cell.text_frame.paragraphs:
                para.line_spacing = 1.0
                for run in para.runs:
                    run.font.name = JP_FONT
                    run.font.size = font_size
                    run.font.color.rgb = INK
                    run.font.bold = c == 0 and r in highlight_rows

    if note:
        box = slide.shapes.add_textbox(
            BODY_LEFT,
            Emu(int(tbl_top) + int(tbl_height) + 60000),
            BODY_WIDTH,
            Emu(400000),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = note
        style_runs(tf, size=Pt(10), color=MUTED)
    return slide


def add_code_slide(prs, title, lead, blocks):
    """blocks: [(見出し, [行, ...])] を等幅風のカードで並べる。"""
    slide = new_content_slide(prs, title, lead)
    drop_shape(slide.placeholders[18])

    n = len(blocks)
    gap = Emu(180000)
    card_h = Emu(int((int(BODY_HEIGHT) * 0.92 - int(gap) * (n - 1)) / n))
    top = BODY_TOP
    for heading, lines in blocks:
        box = slide.shapes.add_textbox(BODY_LEFT, top, BODY_WIDTH, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BG
        box.line.color.rgb = LINE
        box.line.width = Pt(0.75)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(150000)
        tf.margin_right = Emu(120000)
        tf.margin_top = Emu(90000)
        tf.margin_bottom = Emu(90000)
        tf.text = heading
        p0 = tf.paragraphs[0]
        for run in p0.runs:
            run.font.name = JP_FONT
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = ACCENT
        p0.space_after = Pt(3)
        for line in lines:
            para = tf.add_paragraph()
            run = para.add_run()
            run.text = line
            run.font.name = "Consolas"
            run.font.size = Pt(10)
            run.font.color.rgb = INK
            para.space_after = Pt(1)
            para.line_spacing = 1.0
        top = Emu(int(top) + int(card_h) + int(gap))
    return slide


def add_cards_slide(prs, title, lead, cards, *, cols=3, body_size=Pt(11)):
    """cards: [(見出し, [本文行,...])]"""
    slide = new_content_slide(prs, title, lead)
    drop_shape(slide.placeholders[18])

    rows = (len(cards) + cols - 1) // cols
    gap = Emu(200000)
    card_w = Emu(int((int(BODY_WIDTH) - int(gap) * (cols - 1)) / cols))
    card_h = Emu(int((int(BODY_HEIGHT) * 0.88 - int(gap) * (rows - 1)) / rows))

    for i, (heading, lines) in enumerate(cards):
        r, c = divmod(i, cols)
        left = Emu(int(BODY_LEFT) + c * (int(card_w) + int(gap)))
        top = Emu(int(BODY_TOP) + r * (int(card_h) + int(gap)))
        box = slide.shapes.add_textbox(left, top, card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = LINE
        box.line.width = Pt(1.0)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(140000)
        tf.margin_right = Emu(120000)
        tf.margin_top = Emu(110000)
        tf.text = heading
        p0 = tf.paragraphs[0]
        for run in p0.runs:
            run.font.name = JP_FONT
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = ACCENT
        p0.space_after = Pt(6)
        for line in lines:
            para = tf.add_paragraph()
            run = para.add_run()
            run.text = line
            run.font.name = JP_FONT
            run.font.size = body_size
            run.font.color.rgb = INK
            para.space_after = Pt(4)
    return slide


def add_end_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT_END])


def finish(prs, path: Path) -> Path:
    OUTDIR.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path


# --------------------------------------------------------------------------- deck 1
def build_souya_deck() -> Path:
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    add_title_slide(
        prs,
        "CADタグ・属性の自動抽出\n抽出元と具体例のご説明",
        "ICAD／STEP／DXFから何をどう抜き出し、どこまでお渡しできるか",
        "2026.07.28　/　株式会社アルパイン設計事務所　/　創屋様ご説明用",
    )

    add_table_slide(
        prs,
        "本日ご説明する内容",
        "いただいた4つのご質問への回答",
        ["いただいたご質問", "本資料での回答箇所"],
        [
            ["どこからどういったものを抜き出すのかが事前に分かるとありがたい", "CHAPTER 01・02　抽出元カタログと正規化・タグ生成ルール"],
            ["どういったものが抜き出されるのか、具体例が欲しい", "CHAPTER 03　実ICAD 39件の実測結果と実ファイル5例"],
            ["完成後はタグ・属性抽出の部分のソースだけ頂けるとありがたい", "CHAPTER 04　切り出し範囲とモジュール境界"],
            ["ICAD→STEP変換／ナレッジシステム専用ICADライセンスの用意", "CHAPTER 05　変換の実測結果とライセンスの論点"],
        ],
        col_widths=[5, 5],
        font_size=Pt(13),
        header_size=Pt(12),
        height=Emu(2600000),
        note="本資料の事実は、実機ICADでの抽出結果および実装コードに基づいています。未確認の事項は「要確認」と明記しています。",
    )

    # CHAPTER 01
    add_chapter_slide(
        prs, "01", "全体像", "Overview",
        [
            (0, "抽出 → 正規化 → タグ生成 の3段構成"),
            (0, "入力は3経路。ICAD正本／ICADから変換したSTEP・DXF／既存のSTEP・DXF"),
            (0, "抽出器は意味付けをせず、生の値だけを出す"),
        ],
    )

    add_cards_slide(
        prs,
        "処理の流れ",
        "責務を3層に分け、層ごとに切り出せる形にしています",
        [
            ("① 抽出（C#）", [
                "SXNET経由でICADから生の値だけを取得",
                "意味付け・推測は一切しない",
                "1図面 = 1回のプロセス呼び出し",
                "入出力はJSONファイル",
                "出力キー: raw_extract",
            ]),
            ("② 正規化（Python）", [
                "raw_extract を固定スキーマへ整形",
                "2D: 133キー / 3D: 97キー",
                "取れなかった値は null / 空配列のまま残す",
                "「取れなかった」ことも情報として保持",
                "出力キー: canonical_attributes",
            ]),
            ("③ タグ生成（Python）", [
                "canonical_attributes から11種類のタグを生成",
                "辞書はDB化済み。画面から編集可能",
                "各タグに source / evidence / confidence /",
                "reason / tag_rule_version を付与",
                "出力キー: derived_tags",
            ]),
        ],
        cols=3,
        body_size=Pt(12),
    )

    add_table_slide(
        prs,
        "3つの入力経路",
        "ICAD本体が必要になるのは経路Aと、経路Bの変換の瞬間だけ",
        ["経路", "入力", "抽出器", "ICAD本体", "取得できる情報量"],
        [
            ["A", "ICAD .icd 正本", "C# + SXNET", "必要", "最大（材質・質量・部品階層・2D注記まで）"],
            ["B", "ICADから変換したSTEP / DXF", "Python 汎用抽出器", "変換時のみ必要", "中（部品名・階層・レイヤー・文字）"],
            ["C", "客先支給などの既存STEP / DXF", "Python 汎用抽出器", "不要", "中"],
        ],
        col_widths=[1, 3.2, 2.2, 1.6, 4.5],
        font_size=Pt(13),
        height=Emu(2100000),
        note="経路Bと経路Cは同じ抽出器を使います。経路B・Cでは、ICAD本体がなくても抽出処理そのものは動きます（CHAPTER 05のライセンス論点に直結します）。",
    )

    # CHAPTER 02
    add_chapter_slide(
        prs, "02", "抽出元カタログ", "Extraction Sources",
        [
            (0, "ICAD 3D／ICAD 2D／STEP／DXF それぞれの抽出元"),
            (0, "正規化後の属性キー（2D:133 / 3D:97）"),
            (0, "自動タグ11種類と、その生成に使う辞書"),
        ],
    )

    add_table_slide(
        prs,
        "ICAD 3D の抽出元",
        "SXNETに型付きAPIがあり、取得方針が明確な項目",
        ["抽出元", "SXNET根拠", "正規化先（canonical）"],
        [
            ["モデル名・格納フォルダ・コメント", "SxModel.getInf() → SxInfModel.name/path/comment", "model_name, model_path, model_comment"],
            ["パーツ階層", "SxWF.getInfPartTree()", "part_tree_paths"],
            ["パーツ名・コメント", "SxInfPart.name / comment", "part_names, part_comments"],
            ["外部参照・ミラー・未解決参照", "SxInfPart.is_external / is_mirror / is_unloaded", "external_part_exists ほか"],
            ["参照図面名・参照先パス", "SxInfPart.ref_model_name / path", "ref_model_names, ref_model_paths"],
            ["材質（全体・部品単位）", "SxEnt.getInfMaterialList(), SxEntPart.getInfMaterialList()", "material_ids, material_names, material_keywords, part_material_candidates"],
            ["質量・重量・体積・面積・密度", "SxEnt.getMass() → SxInfMass", "mass_value, weight_value, volume_value, area_value, density_value"],
            ["重心・慣性モーメント", "SxInfMass.pos, inf_global_moment ほか", "center_of_gravity, global_moment, inertia_moment_candidates"],
            ["トップパーツ付加情報", "SxWF.getInfExTopPart(), SxInfPartTree.ex_inf", "top_part_ex_info, part_ex_info_fields"],
        ],
        col_widths=[3.2, 4.3, 4.5],
        font_size=Pt(10.5),
        header_size=Pt(11),
    )

    add_table_slide(
        prs,
        "ICAD 2D の抽出元",
        "図枠は文字・注記から候補化し、座標はレビュー証跡として保持する",
        ["抽出元", "raw_extract", "正規化先（canonical）"],
        [
            ["ビューシート名・尺度・種別", "view_sheets[]", "model_view_sheet_count, scale_candidates"],
            ["出図範囲枠・用紙サイズ", "print_frames[]", "paper_size, 印刷枠内外判定"],
            ["文字・注記", "texts[]", "text_tokens, label_texts"],
            ["図枠のラベルと値", "texts[]（同一要素の同じ行／次行）", "title_block_candidates → title_block_fields"],
            ["寸法", "dimensions[]", "dimension_values, dimension_symbols"],
            ["公差・幾何公差", "tolerances[]", "tolerance_candidates"],
            ["溶接記号・注記", "weld_notes[]", "weld_note_candidates"],
            ["バルーン", "balloons[]", "balloon_candidates"],
            ["面粗さ・仕上げ記号", "geometry_primitives[]", "surface_roughness_values, finish_mark_types"],
            ["切断線・矢視・ハッチング", "geometry_primitives[]", "curve_section_candidates, cut_line_count, hatch_or_section_count"],
            ["穴・長穴候補", "geometry_primitives[]", "hole_candidate_count, slot_candidate_count"],
        ],
        col_widths=[3.2, 3.6, 5.2],
        font_size=Pt(10),
        header_size=Pt(11),
        note="図番・図面名・担当者・承認者・日付・材質・重量・表面処理・塗装指示・PRFX・ユニット番号は、SXNETの固定フィールドとしては存在しません。会社ごとの図枠差が大きいため、title_block_fields は固定列ではなく key-value 候補として持ち、採用条件を通ったものだけを確定値に上げています。",
    )

    add_table_slide(
        prs,
        "STEP / DXF の抽出元",
        "外部CADライブラリを使わず、標準ライブラリだけで解析しています",
        ["形式", "抽出元", "raw_extract", "正規化先（canonical）"],
        [
            ["STEP", "ヘッダ（FILE_NAME, FILE_DESCRIPTION）", "model_info.name / comment", "model_name, model_comment"],
            ["STEP", "PRODUCT / PRODUCT_DEFINITION", "step_products[]", "step_product_names"],
            ["STEP", "NEXT_ASSEMBLY_USAGE_OCCURRENCE", "step_assembly_relationships[]", "step_assembly_relationship_count, parts[].tree_path"],
            ["STEP", "文字列中の材質パターン", "materials[]", "material_keywords"],
            ["DXF", "TEXT / MTEXT / ATTRIB / ATTDEF", "texts[]", "text_tokens, title_block_candidates"],
            ["DXF", "INSERT + ATTRIB のブロック参照", "block_references[].attributes[]", "dxf_block_references, dxf_block_attribute_tokens"],
            ["DXF", "DIMENSION の表示値候補", "dimensions[]", "dimension_values"],
            ["DXF", "LINE / CIRCLE / ARC / LWPOLYLINE / HATCH ほか", "geometry_primitives[]", "形状特徴候補"],
            ["DXF", "レイヤー一覧", "layers[]", "dxf_layers"],
        ],
        col_widths=[1, 4, 3.6, 4],
        font_size=Pt(10.5),
        header_size=Pt(11),
        note="両形式ともファイルパス・ファイル名を source_path_tokens として保持し、客先・案件・装置カテゴリの辞書照合に使います。",
    )

    add_bullet_slide(
        prs,
        "正規化後の属性キー",
        "固定スキーマ。2D snapshot で133キー、3D snapshot で97キー",
        [
            (0, "図面識別　drawing_number, drawing_name, revision, drawing_size, paper_size, scale, document_kind"),
            (0, "担当・日付　designer, checker, approver, drawing_date, created_date, checked_date, approved_date, revision_date"),
            (0, "分類　customer_name, project_name, equipment_name, equipment_category, module_name, prfx, unit_number, owner, status"),
            (0, "ファイル出所　source_full_path, source_directory_path, source_file_name, source_file_stem, source_path_tokens"),
            (0, "質量特性（3D）　mass_value, weight_value, volume_value, area_value, density_value, center_of_gravity, global_moment, inertia_moment_candidates"),
            (0, "材質　material_ids, material_names, material_keywords, unresolved_material_keywords, part_material_candidates"),
            (0, "部品構成　part_names, part_comments, part_tree_paths, ref_model_names, external_part_exists, mirror_part_exists, unresolved_part_exists"),
            (0, "2D注記・図枠　title_block_fields, title_block_candidates, revision_note_candidates, text_tokens, raw_2d_sections"),
            (0, "2D寸法・記号　dimension_values, tolerance_candidates, weld_note_candidates, balloon_candidates, surface_roughness_values, hole_candidate_count, slot_candidate_count"),
            (0, "加工・処理　surface_treatment_tokens, paint, heat_treatment_keywords, hardness_spec_values, process_keywords"),
            (0, "キーワード　part_keywords, maker_keywords, spec_tokens, inspection_keywords, change_keywords, issue_keywords"),
            (0, "STEP/DXF固有　step_products, step_product_names, step_assembly_relationships, dxf_layers, dxf_block_references"),
            (0, "抽出メタ　extraction_status, ocr_used, confidence_summary, normalizer_version"),
            (1, "値が取れなかったキーは null / 空配列のまま残します。「取れなかった」ことも情報として扱う方針です。"),
        ],
        base_size=Pt(12),
    )

    add_table_slide(
        prs,
        "自動タグ11種類と辞書",
        "各タグに source / evidence / confidence / reason を付与。根拠なしのタグは作らない",
        ["タグ接頭辞", "生成元の属性キー", "確度", "辞書エントリ数（初期値）"],
        [
            ["客先:", "customer_name", "high", "3（コマツ小山、広島アルミ、澁谷工業）"],
            ["案件:", "project_name", "high", "初期0件（画面から登録）"],
            ["装置:", "equipment_category", "high", "30（ガントリー、治具、コンベア、架台 ほか）"],
            ["メーカー:", "maker_keywords", "medium", "40（SMC、ミスミ、THK、NSK ほか）"],
            ["材質:", "material_keywords, title_block_fields.material", "medium", "68（SS400、SUS304、S45C、FC300 ほか）"],
            ["表面処理:", "surface_treatment_tokens", "medium", "―"],
            ["塗装:", "paint_instruction_tokens", "medium", "―"],
            ["熱処理:", "heat_treatment_keywords", "medium", "20（焼入れ、調質、浸炭、窒化 ほか）"],
            ["PRFX:", "prfx_candidates, title_block_fields.prfx", "medium", "―"],
            ["ユニット:", "unit_number_candidates", "medium", "―"],
            ["規格:", "spec_tokens のうち辞書一致語（現状 SES のみ）", "medium", "7（SES、JIS、ISO、DIN、ANSI ほか）"],
        ],
        col_widths=[1.6, 4.6, 1.2, 4.6],
        font_size=Pt(10.5),
        header_size=Pt(11),
        note="辞書はDB化済みで、画面（システム設定 > タグ辞書管理）と /admin から編集できます。案件辞書は初期0件で、画面からの登録が正本です。照合対象はファイルパストークンだけでなく、モデル情報・図面内文字・DXFレイヤー等を合成した検索トークン列（part_keywords）です。なお spec_tokens は2Dでは図面内文字と公差テキストの生の集合（実例で1図面112件）で、その中の辞書一致語だけを 規格: タグにしています。",
    )

    add_cards_slide(
        prs,
        "タグ化するもの／属性保持に留めるもの",
        "存在するだけでタグにすると検索ノイズが大きくなるため、意図的に分けています",
        [
            ("自動タグ化する", [
                "客先 / 案件 / 装置カテゴリ",
                "メーカー名",
                "正式材質",
                "表面処理・塗装",
                "熱処理",
                "PRFX・ユニット番号",
                "SES などの明確な規格識別子",
            ]),
            ("属性として保持し、タグ化しない", [
                "寸法値・公差値",
                "溶接記号・バルーン",
                "穴・長穴・切断線・ハッチング等の形状特徴",
                "質量・体積・重心・慣性モーメント",
                "",
                "→ 図面レビューやRAG投入時の",
                "　 属性・根拠として保持します",
            ]),
        ],
        cols=2,
        body_size=Pt(13),
    )

    # CHAPTER 03
    add_chapter_slide(
        prs, "03", "具体例", "Real Extraction Results",
        [
            (0, "共有いただいた実ICAD 39件の実測結果"),
            (0, "実ファイル5例の、実際の出力値（値は加工していません）"),
            (0, "具体例から読み取れる特性と、すり合わせが必要な点"),
        ],
    )

    add_table_slide(
        prs,
        "実ICAD 39件の実測結果",
        "共有いただいた実ICAD 39件を、抽出から画面表示まで通した結果（2026-07-17時点）",
        ["観点", "結果"],
        [
            ["抽出完了", "39件すべてで 2D snapshot / 3D snapshot の両方を保存済み、未抽出0件"],
            ["部品名（part_names）", "39 / 39 件で取得"],
            ["質量・重量（mass_value / weight_value）", "38 / 39 件で取得"],
            ["材質（material_keywords）", "33 / 39 件で取得"],
            ["生成タグ（上位）", "材質:SS400 13件　/　材質:SUS 13件　/　材質:SUS304 12件　/　メーカー:SMC 4件　/　材質:ねずみ鋳鉄 4件"],
            ["", "材質:SUS304B・SUS316・SPCC・A5052P・S45C・FC300 各3件　/　規格:SES 2件　/　客先:澁谷工業 1件　/　装置:治具 1件"],
        ],
        col_widths=[3.4, 8],
        font_size=Pt(12),
        height=Emu(2700000),
        note="材質タグは安定して出ます。客先・装置タグが少ないのは、この測定時点の客先辞書が3件しか入っていなかったためです（辞書のDB化は2026-07-17）。辞書を実運用の語彙で拡充すれば増える性質のもので、抽出ロジックの限界ではありません。",
    )

    add_code_slide(
        prs,
        "具体例①　客先・装置がパスから確定したケース",
        "実ICADファイルの実際の出力（値は未加工）",
        [
            ("U8105111315.icd", [
                r"パス   : J:\SBY\アイソレータ\210126_エーザイ_アイソレータ_RAA4844\作業フォルダ\開閉扉\U8105111315.icd",
                "タグ   : 客先:澁谷工業 / 材質:SUS304 / 材質:SUS / 規格:SES",
                ' 属性  : customer_name="澁谷工業", mass_value=0.18021418, weight_value=1.7672974,',
                '         material_keywords=["SUS304","SUS"], part_names=["U81051113150"]',
            ]),
            ("XH30-A08001-R03-JP_ロードカップ部改造.icd", [
                r"パス   : J:\ZCSET\300P_210312\作業\2_ロードカップ部\XH30-A08001-R03-JP_ロードカップ部改造.icd",
                "タグ   : 装置:治具 / メーカー:SMC / 材質:SUS316 / 材質:PVC / 材質:PPS / 材質:PTFE /",
                "         材質:A5052P / 材質:SUS304 / 材質:SUS / 材質:NBR / 材質:POM",
                ' 属性  : equipment_category="治具", mass_value=4.75944799, weight_value=46.67424068,',
                '         part_names 108件（"＠リフト部", "CSP300R-8012-00_ノズルブラケット-1", "Oリング(S190)" ほか）',
                '         unresolved_material_keywords=["75"]  ← 材質と断定できなかった値は別枠に隔離',
            ]),
        ],
    )

    add_code_slide(
        prs,
        "具体例②　大規模アセンブリ・規格・塗装",
        "実ICADファイルの実際の出力（値は未加工）",
        [
            ("474300AC219.icd　（部品189件の大規模アセンブリ）", [
                r"パス   : \\HONSYA-FILE01\data_cad3d\SBY\CAP\260527_AAM6351_アイリスオーヤマ_宮本様\474300AC219.icd",
                "タグ   : メーカー:SMC / 材質:SUS304 / 材質:PET / 材質:H-PVC / 材質:NBR / 材質:EPDM ほか計19タグ",
                " 属性  : mass_value=17.7113085, weight_value=173.68860346, material_keywords 18種, part_names 189件",
                '         unresolved_material_keywords=["ZZZ"]',
            ]),
            ("TR1D9K99027.icd　（規格・面粗さ・公差）", [
                r"パス   : J:\シブヤパッケージングシステム\25_9R_膨潤パレットアキューム部\...\部品図(新規、訂正)\9K\TR1D9K99027.icd",
                "タグ   : 材質:A5000 / 材質:SUS304 / 材質:A1000 / 材質:SUS / 規格:SES",
                " 属性  : mass_value=2.15703985, spec_tokens 112件（\"Ra 6.3\", \"粗級\", \"中級\", \"±4\", \"±8\", \"±2\" ほか）",
            ]),
            ("03_20K03379P00_ｼｭｰﾄﾍﾞｰｽ(No.2FFS_XS).icd　（図枠から塗装指示）", [
                r"パス   : J:\アースエンジニアリング\251216_ツネイシカムテックス\...\AR05-A05-B04_No.2 F-FスクリーンスクリーンSシュート組立図\...",
                "タグ   : 材質:SS400 / 塗装:ﾊ仕様書ﾆﾖﾙ",
                ' 属性  : mass_value=11.66198417, part_names=["03_20K03379P00_...", "溝形鋼_100*50*5*7.5"]',
            ]),
        ],
    )

    add_bullet_slide(
        prs,
        "具体例から読み取れる特性",
        "実データを通したうえでの所見と、すり合わせが必要な点",
        [
            (0, "パスは強力な情報源です", True),
            (1, "客先・案件・装置は、図面内の文字よりフォルダ構成から確定できるケースが多い"),
            (1, "source_path_tokens を含む検索トークン列（part_keywords）を辞書照合しています"),
            (0, "材質は3Dの材質APIから安定して取れます", True),
            (1, "ただし ZZZ / CDQ / 75 のような、材質と断定できない値が混ざります"),
            (1, "これらは unresolved_material_keywords に隔離し、タグにはしません"),
            (0, "質量・重量は3Dからほぼ確実に取れます（38/39）", True),
            (0, "半角カナ・機種依存文字がそのまま入ります", True),
            (1, "実例: ｼｭｰﾄﾍﾞｰｽ、ﾊ仕様書ﾆﾖﾙ"),
            (1, "表示・検索時の正規化方針（半角カナの扱い、全角統一の要否）はすり合わせが必要です"),
        ],
        base_size=Pt(16),
    )

    # CHAPTER 04
    add_chapter_slide(
        prs, "04", "ソースの\n切り出し範囲", "Module Boundary",
        [
            (0, "「タグ・属性抽出の部分だけ切り出したい」というご要望に対する回答"),
            (0, "C#側は独立ソリューション。ナレッジシステム本体への依存なし"),
            (0, "Python側のDjango結合は実質2点だけ"),
        ],
    )

    add_table_slide(
        prs,
        "C#側（ICAD抽出コア）",
        "独立ソリューション（IcadExtraction.sln、約4,700行）としてそのままお渡しできます",
        ["プロジェクト", "役割"],
        [
            ["IcadExtraction.Contracts", "入出力のデータ契約（JSONスキーマ相当）"],
            ["IcadExtraction.SxNet", "SXNET経由の2D/3D抽出、材質・質量プローブ、STL/STEP/DXF出力"],
            ["IcadExtraction.Runner", "CLIエントリ。extract / extract-batch / detect / probe-2d-print / convert-cad /\nprobe-cad-export-types / cancel / clear-command / shutdown-icad / self-check の10コマンド"],
        ],
        col_widths=[3, 8],
        font_size=Pt(13),
        height=Emu(1900000),
        note="インターフェースは「1図面 = 1回のプロセス呼び出し」、入出力はJSONファイルです。呼び出し側の言語を選びません。ナレッジシステム本体への依存はありません。",
    )

    add_table_slide(
        prs,
        "Python側（正規化・タグ生成コア）",
        "この6ファイル（約2,800行）が本体です",
        ["ファイル", "行数", "役割", "外部依存"],
        [
            ["services/normalization.py", "1,860", "raw_extract → canonical_attributes", "settings 1定数、TagDictionaryEntry（定数参照のみ）"],
            ["services/generic_cad_extractor.py", "519", "STEP/DXFの抽出（外部CADライブラリ不要）", "settings 1定数"],
            ["services/seed_dictionaries.py", "222", "初期辞書（純Python）", "なし"],
            ["services/tag_builder.py", "112", "canonical_attributes → derived_tags", "settings 1定数"],
            ["services/dictionaries.py", "54", "辞書ロード（DB + 初期辞書のマージ）", "TagDictionaryEntry（ORMクエリあり）"],
            ["services/source_formats.py", "41", "拡張子→フォーマット判定（純Python）", "なし"],
        ],
        col_widths=[3.6, 1, 3.6, 3.2],
        font_size=Pt(11.5),
        header_size=Pt(11),
        height=Emu(2500000),
        note="Djangoへの結合は実質2点だけです。① settings のバージョン文字列3個（NORMALIZER_VERSION / TAG_RULE_VERSION / SCHEMA_VERSION）、② TagDictionaryEntry モデル（辞書のDB読み出し）。①は引数注入で置き換えるだけです。②はDBアクセスを伴いますが、dictionaries.load_keyword_mapping の1箇所に集約されており、辞書ロード関数をインターフェース化すれば分離できます。",
    )

    add_cards_slide(
        prs,
        "お渡しできる形と、含めないもの",
        "ご希望の形に合わせます",
        [
            ("ご希望であればこの形で納品します", [
                "icad_tag_extraction 単体パッケージ（Django非依存）",
                "サンプル入出力JSON",
                "単体テスト",
                "辞書の初期データ",
                "スキーマ定義",
                "",
                "C#側は IcadExtraction.sln 一式",
            ]),
            ("切り出し対象に含めないもの", [
                "Djangoモデル（RegisteredDrawing,",
                "　DrawingMetadataSnapshot, ExtractionJob,",
                "　AuditLog, TagDictionaryEntry）",
                "ジョブ管理・リトライ・タイムアウト制御・監査ログ",
                "画面（タグ辞書管理、抽出管理、レビューUI）",
                "",
                "→ 参考実装としてはお渡しできますが、",
                "　 そのまま組み込む前提では作っていません",
            ]),
        ],
        cols=2,
        body_size=Pt(12.5),
    )

    # CHAPTER 05
    add_chapter_slide(
        prs, "05", "STEP変換と\nICADライセンス", "Conversion & License",
        [
            (0, "ICAD→STEP/DXF変換は実機で確認済み。ただし限界がある"),
            (0, "「ナレッジシステム専用ICAD」が技術的に必須という結論にはならない"),
            (0, "販売を前提にした場合の選択肢を整理"),
        ],
    )

    add_cards_slide(
        prs,
        "ICAD→STEP / DXF 変換の実測と限界",
        "2026-07-26 に実機で確認した内容",
        [
            ("できたこと", [
                "実機SXNETの出力形式定数を確認",
                "　FILE_TYPE_STEP=11 / FILE_TYPE_DXF=1",
                "　形式別の数値オーバーライドなしで変換可能",
                "9NK452WX90-00-LINER-A3-3D-01.icd を",
                "STEP / DXF へ変換し、抽出・snapshot保存まで完了",
                "変換後STEP: step_product_names=",
                '　["Assembly", "...-prt0"], 親子関係1件',
                "変換後DXF: dxf_layers 5件を取得",
            ]),
            ("限界（重要）", [
                "ICAD本体にある材質・質量は、STEP側に同等には残らない",
                "　→ 変換後データはICAD正本と等価ではない",
                "今回のDXFサンプルではブロック属性は0件",
                "　→ 図枠情報の有無は客先の図枠仕様に依存",
                "STEPはSXNETが .stp で出力（.step/.stp 両対応済み）",
                "DXF変換時、export成功後もrunner終了待ちが長いケースあり",
                "変換後の保存確認ダイアログは保存不要",
                "　（shutdown-icad で保存なし終了）",
            ]),
        ],
        cols=2,
        body_size=Pt(11.5),
    )

    add_bullet_slide(
        prs,
        "変換に関する所見",
        "変換経由は「上位互換」ではなく「代替経路」です",
        [
            (0, "「ICAD→STEP変換を機能として用意する」ことは技術的に成立します", True),
            (0, "ただし、変換したSTEPから取れる情報は、ICAD正本から直接抜いた情報より確実に少ないという実測結果があります", True),
            (1, "材質・質量が落ちるのが最も大きい差です"),
            (0, "変換を経由するのは「ICADライセンスのない環境でも最低限の情報を取りたい」場合の代替経路であり、上位互換ではありません", True),
            (0, "したがって、製品仕様としては「経路Aで取れる情報」と「経路B・Cで取れる情報」の差を明示し、顧客の環境に応じて構成を選べる形にすることを提案します", True),
        ],
        base_size=Pt(16),
    )

    add_table_slide(
        prs,
        "ICADライセンス｜技術的な事実",
        "コード上で確認済みの事実。ここから判断の前提を揃えます",
        ["#", "事実", "根拠"],
        [
            ["1", "経路A（ICAD正本から抽出）は、SXNET経由でICAD SX本体プロセスを必要とする", "IcadProcessStarter がICAD実行体を起動、または起動済みプロセスへ接続"],
            ["2", "ICADセッションは1つずつ排他利用される。並列に複数の抽出は走らせない", "Local\\KnowledgeSystem.IcadExtraction.IcadSession の Mutex で直列化（排他範囲はWindowsのログオンセッション単位）"],
            ["3", "ICADが必要なのは登録・抽出のタイミングだけ。検索・閲覧・タグ表示・RAG回答にICADは不要", "抽出結果はsnapshotとしてDBに保存され、以降は参照されない"],
            ["4", "経路B・C（STEP/DXFからの抽出）はICADなしで動く。外部CADライブラリも不要", "generic_cad_extractor.py は外部CADライブラリに依存せず、Python標準の文字列処理だけで解析"],
            ["5", "ICAD→STEP/DXF変換の、その瞬間だけICADが要る", "IcadCadFormatExporter（SXNET SxModel.export）"],
        ],
        col_widths=[0.5, 5.5, 5.5],
        font_size=Pt(11),
        header_size=Pt(11),
        height=Emu(3000000),
        note="したがって「ナレッジシステム専用のICADライセンス」が技術的に必須という結論にはなりません。必要なのは「抽出処理を実行している間、そのICADセッションを占有できること」だけです。",
    )

    add_table_slide(
        prs,
        "販売を前提にした場合の選択肢",
        "「顧客が150万円のICADを追加購入する」を前提にした製品は売りにくい。前提を分けて設計する",
        ["案", "内容", "顧客の追加ICAD費用", "取得できる情報", "主な制約"],
        [
            ["① 既存ライセンス兼用\n（夜間バッチ）", "顧客が既に持つICADを、設計者が使わない夜間・休日に抽出用として使う", "0円", "最大", "ライセンス条項の可否確認が必須。抽出中はその座席を占有"],
            ["② 既存ライセンス兼用\n（専用PC1台）", "顧客の空き座席1本を抽出専用PCへ割り当て", "0円\n（座席の振替）", "最大", "空き座席がある顧客に限る"],
            ["③ ICAD不要構成", "顧客側でSTEP/DXF/PDFへ出力済みのデータを取り込む（経路B・Cのみ）", "0円", "中\n（材質・質量は落ちる）", "顧客に出力運用を依頼する必要あり"],
            ["④ 変換代行", "当社の余剰ライセンス＋当社サーバーでICAD→STEP/中間データ変換を代行", "0円", "最大〜中", "図面データを社外（当社）へ出すことへの顧客同意が必要"],
            ["⑤ 専用ライセンス\n新規購入", "顧客がナレッジシステム用にICADを1本追加", "約150万円", "最大", "費用が導入判断のボトルネックになりやすい"],
        ],
        col_widths=[1.6, 4, 1.5, 1.4, 3],
        font_size=Pt(10),
        header_size=Pt(10.5),
        highlight_rows={1, 4},
        note="当社には余剰ライセンスがあるため、④は当社が担げる現実的な選択肢です。また対象顧客はICADユーザーであることが多く、①②が成立する可能性も相応にあると見ています。",
    )

    add_bullet_slide(
        prs,
        "要確認事項とご相談",
        "ライセンス条項の問題であり、当社もまだ確認できていません。断定していません",
        [
            (0, "当社 → 富士通／ICAD販売元 へ確認する事項", True),
            (1, "設計者向けライセンスを、無人のバッチ処理（サーバー常駐・自動起動）で使うことが許諾範囲か"),
            (1, "ライセンス形態（ノードロック／フローティング等）と、フローティングの場合の借用可否"),
            (1, "「同時使用」の定義。1座席を夜間だけ別PCで使う運用が可能か"),
            (1, "当社の余剰ライセンスを使って他社図面の変換を受託することが許諾範囲か（案④の前提）"),
            (0, "創屋様へのご相談", True),
            (1, "製品としては「ICADがある構成」と「ICADがない構成」の両方を成立させる方向で進めたいと考えています"),
            (1, "経路Aで取れる情報と経路B・Cで取れる情報の差を仕様として明示し、顧客環境に応じて構成を選べる形です"),
            (1, "この方針で問題ないか、ご意見をいただきたいところです"),
        ],
        base_size=Pt(15),
    )

    add_table_slide(
        prs,
        "創屋様への確認事項",
        "本資料に対するご回答をいただきたい項目",
        ["区分", "確認事項"],
        [
            ["抽出", "1. STEPから製品名・部品名・部品階層・材質・質量特性を取得できるライブラリ／APIを、創屋様側でお持ちか"],
            ["抽出", "2. DXFからTEXT/MTEXT、ブロック属性、レイヤー名、寸法、公差、溶接記号を分離取得できるか"],
            ["抽出", "3. 図枠のラベルと値を同一TEXT/MTEXT要素またはDXFブロック属性として取得できるか（別要素の座標ペアリングは対象外）"],
            ["抽出", "4. 材質が色・レイヤー・ブロック名にしか入っていないケースがあるか"],
            ["抽出", "5. 抽出値に推測を混ぜず、取得元フィールドと信頼度を添えて返せるか"],
            ["受け渡し", "6. Django非依存の純Pythonパッケージ形での納品でよいか。それとも現状のDjango app のまま渡す方が組み込みやすいか"],
            ["受け渡し", "7. 一緒に欲しい成果物（サンプル入出力JSON、単体テスト、辞書初期データ、スキーマ定義）の優先順位"],
            ["受け渡し", "8. タグ・属性の書き込み先API（drawing_attributes / product_attributes / part_attributes 相当）の確定仕様"],
            ["構成", "9. 「ICADあり／なし両構成」の方針で問題ないか"],
            ["構成", "10. ICAD→STEP変換を創屋様側で実装される場合、当社のC# convert-cad 実装をそのまま使うか、作り直すか"],
        ],
        col_widths=[1.2, 10],
        font_size=Pt(11),
        header_size=Pt(11),
    )

    add_end_slide(prs)
    return finish(prs, OUTDIR / "20260728_ナレッジシステム_タグ属性抽出_創屋様向けご説明_r1.pptx")


# --------------------------------------------------------------------------- deck 2
def build_exec_deck() -> Path:
    prs = Presentation(TEMPLATE)
    clear_slides(prs)

    add_title_slide(
        prs,
        "ナレッジシステム\nCADタグ・属性の自動抽出",
        "開発到達点・業務インパクト・ICADライセンスのご判断事項",
        "2026.07.28　/　株式会社アルパイン設計事務所　/　役員会ご説明用",
    )

    add_cards_slide(
        prs,
        "要旨",
        "本日ご判断いただきたいことは3点です",
        [
            ("① 開発は動く段階まで来ている", [
                "実ICAD 39件で、抽出→正規化→タグ生成→",
                "画面表示まで一気通貫で動作を確認済み。",
                "",
                "材質・質量・部品構成は実データで安定取得。",
                "",
                "残るのは辞書の拡充と、",
                "創屋への受け渡し形式の確定。",
            ]),
            ("② ICAD専用ライセンスは必須ではない", [
                "「ナレッジシステム専用に150万円の",
                "ICADを追加購入」が前提だと売りにくい。",
                "",
                "技術的には、ICADが要るのは",
                "登録・抽出の瞬間だけ。",
                "",
                "既存ライセンス兼用／当社の余剰ライセンス",
                "による変換代行／ICAD不要構成が取り得る。",
            ]),
            ("③ 外販を見据えた構成判断が要る", [
                "「ICADがある構成」と「ICADがない構成」の",
                "両方を成立させる方向で進めたい。",
                "",
                "ライセンス条項の可否は富士通・販売元へ",
                "確認が必要（未確認）。",
                "",
                "→ この方針でよいかのご判断をお願いします。",
            ]),
        ],
        cols=3,
        body_size=Pt(11.5),
    )

    # CHAPTER 01
    add_chapter_slide(
        prs, "01", "現在の開発到達点", "Progress",
        [
            (0, "何ができるようになったか"),
            (0, "実データでどこまで取れたか"),
            (0, "図面1枚から実際に何が出るか"),
        ],
    )

    add_cards_slide(
        prs,
        "何ができるようになったか",
        "図面ファイルを登録すると、人手を介さずタグと属性が付く",
        [
            ("① CADから中身を読み出す", [
                "ICADの2D図面・3Dモデルから、",
                "材質・質量・部品構成・寸法・注記・",
                "図枠情報を機械的に取り出す。",
                "",
                "ICADのAPI（SXNET）を直接使用。",
                "画面操作の自動化ではない。",
            ]),
            ("② 意味のある属性に整える", [
                "取り出した生データを、検索に使える",
                "固定項目へ整形する。",
                "",
                "2D図面で133項目、3Dモデルで97項目。",
                "",
                "取れなかった項目は「取れなかった」と",
                "記録し、推測で埋めない。",
            ]),
            ("③ タグを自動で付ける", [
                "客先・案件・装置・メーカー・材質・",
                "表面処理・塗装・熱処理・PRFX・",
                "ユニット・規格の11種類。",
                "",
                "全タグに根拠と確度が付く。",
                "根拠のないタグは作らない。",
            ]),
        ],
        cols=3,
        body_size=Pt(12),
    )

    add_table_slide(
        prs,
        "実データでの実績",
        "共有した実ICAD 39件を、実際に最後まで通した結果（2026-07-17時点）",
        ["観点", "結果", "評価"],
        [
            ["抽出完了", "39 / 39 件（未抽出0件）", "◎ 落ちずに通る"],
            ["部品構成（部品名）", "39 / 39 件で取得", "◎ 全件取得"],
            ["質量・重量", "38 / 39 件で取得", "◎ ほぼ全件"],
            ["材質", "33 / 39 件で取得", "○ 8割超"],
            ["客先・装置のタグ", "少数（辞書が検証用に3件のみだったため）", "△ 辞書拡充で改善する見込み"],
        ],
        col_widths=[2.8, 5.5, 3],
        font_size=Pt(13),
        header_size=Pt(12),
        height=Emu(2600000),
        highlight_rows={5},
        note="客先・装置タグが少ないのは、測定時点で客先辞書に3件しか登録していなかったためで、抽出ロジックの限界ではありません。辞書はDB化済みで画面から追加できます。運用時に実際の客先・装置語彙を登録する作業が必要です。",
    )

    add_code_slide(
        prs,
        "図面1枚から実際に何が出るか",
        "実ICADファイルの実際の出力（値は未加工）",
        [
            ("例1　澁谷工業向けの部品図（U8105111315.icd）", [
                "タグ  : 客先:澁谷工業 / 材質:SUS304 / 材質:SUS / 規格:SES",
                "属性  : 質量 0.18kg（mass_unit_name=mm-kg）、部品名 U81051113150",
                "→ 客先×材質での絞り込みに使えるデータが、人手を介さず揃う",
            ]),
            ("例2　治具の3Dモデル（XH30-A08001-R03-JP_ロードカップ部改造.icd）", [
                "タグ  : 装置:治具 / メーカー:SMC / 材質:SUS316, PVC, PPS, PTFE, A5052P, SUS304, NBR, POM",
                "属性  : 質量 4.76kg、構成部品 108点",
                "→ 「SMCを使った治具」「PTFEを使った案件」を横断で絞り込める土台になる",
            ]),
            ("例3　大規模アセンブリ（474300AC219.icd）", [
                "タグ  : メーカー:SMC / 材質18種 ほか計19タグ",
                "属性  : 質量 17.7kg、構成部品 189点",
                "→ 189点の部品を人手で棚卸しせずに、材質構成をデータとして持てる",
            ]),
        ],
    )

    # CHAPTER 02
    add_chapter_slide(
        prs, "02", "業務インパクト", "Business Impact",
        [
            (0, "設計者の業務が何に変わるか"),
            (0, "効果が出るための前提条件"),
            (0, "現時点で定量化できていないこと"),
        ],
    )

    add_cards_slide(
        prs,
        "何が変わるか",
        "「探す時間」を「設計する時間」に振り替える",
        [
            ("設計者の探索時間", [
                "現状：過去案件の図面を探すのに、",
                "フォルダを辿る／記憶に頼る／人に聞く。",
                "",
                "導入後：客先・装置・材質・メーカーで",
                "横断検索し、根拠付きで候補が出る。",
            ]),
            ("図面資産の再利用", [
                "現状：類似案件があっても、",
                "存在を知らなければ使えない。",
                "",
                "導入後：材質・部品構成・メーカーの",
                "組み合わせから類似図面に到達できる。",
            ]),
            ("属人性の低減", [
                "現状：どこに何があるかがベテランの",
                "記憶に依存している。",
                "",
                "導入後：タグと属性がデータとして残り、",
                "経験年数に依らず同じ探索ができる。",
            ]),
            ("調達・見積の根拠", [
                "現状：材質構成・質量は図面を開いて確認。",
                "",
                "導入後：質量・材質・部品点数が",
                "属性として蓄積され、集計できる。",
            ]),
            ("設計品質", [
                "現状：過去のトラブル事例が",
                "個人のメモに閉じている。",
                "",
                "導入後：図面とタグを起点に、",
                "関連資料へ辿れる土台ができる。",
            ]),
            ("外販の可能性", [
                "同じ課題はICADを使う設計会社に共通。",
                "",
                "自社で運用実績を作ったうえで、",
                "製品として外部提供する道がある。",
            ]),
        ],
        cols=3,
        body_size=Pt(11),
    )

    add_bullet_slide(
        prs,
        "効果の前提条件と、まだ言えないこと",
        "過大評価を避けるため、前提と未計測を明示します",
        [
            (0, "効果を出すために必要な作業（技術ではなく運用の話）", True),
            (1, "辞書の拡充：客先・装置・案件の語彙を実運用ベースで登録する"),
            (1, "フォルダ命名規則の整備：パスから客先・案件を確定できることが精度に直結する"),
            (1, "登録運用の定義：どのタイミングで、誰が、どの図面を登録するか"),
            (1, "レビュー運用：自動付与タグを人が確認・補正する導線"),
            (0, "現時点で定量化できていないこと", True),
            (1, "「検索時間が何割減るか」は未計測です。実運用データがないため、数値は出せません"),
            (1, "本資料に時間削減率・金額効果の数値を載せていないのは、根拠がないためです"),
            (1, "定量化するには、社内で試験運用を行い、導入前後の実測を取る必要があります"),
            (0, "→ ご提案：まず自社で試験運用し、定量効果を測ってから外販判断へ進む", True),
        ],
        base_size=Pt(15),
    )

    # CHAPTER 03
    add_chapter_slide(
        prs, "03", "ICADライセンス\nのご判断", "License Decision",
        [
            (0, "創屋から「専用ICAD 約150万円を販売先が用意してくれるか」という懸念が出ている"),
            (0, "技術的な事実を踏まえると、専用ライセンスは必須ではない"),
            (0, "販売を前提とした場合の選択肢を整理した"),
        ],
    )

    add_bullet_slide(
        prs,
        "論点の整理",
        "創屋からの問いかけと、当社が確認した事実",
        [
            (0, "創屋からの問いかけ", True),
            (1, "「ナレッジシステム専用として、150万円程のICADを販売先が用意してくれるかが当初からの懸念」"),
            (1, "「ナレッジシステム専用である必要はあるのか。兼用は難しいか」"),
            (0, "当社がコード上で確認した事実", True),
            (1, "ICADが必要なのは、図面を登録して情報を抜き出すタイミングだけ"),
            (1, "検索・閲覧・タグ表示・RAG回答にICADは不要（抽出結果はDBに保存済みのため）"),
            (1, "抽出処理はICADを1セッションずつ占有する。複数を並列に走らせる作りではない"),
            (1, "STEP/DXFからの抽出はICADなしで動く（ただし材質・質量は落ちる）"),
            (0, "結論", True),
            (1, "「ナレッジシステム専用のICADライセンス」が技術的に必須という結論にはならない"),
            (1, "必要なのは「抽出処理を実行している間、そのICADセッションを占有できること」だけ"),
        ],
        base_size=Pt(15),
    )

    add_table_slide(
        prs,
        "販売を前提にした場合の選択肢",
        "「顧客が150万円のICADを追加購入する」を前提にした製品は売りにくい",
        ["案", "内容", "顧客の追加ICAD費用", "取得できる情報", "主な制約"],
        [
            ["① 既存ライセンス兼用\n（夜間バッチ）", "顧客が既に持つICADを、設計者が使わない夜間・休日に抽出用として使う", "0円", "最大", "ライセンス条項の可否確認が必須。抽出中はその座席を占有"],
            ["② 既存ライセンス兼用\n（専用PC1台）", "顧客の空き座席1本を抽出専用PCへ割り当て", "0円\n（座席の振替）", "最大", "空き座席がある顧客に限る"],
            ["③ ICAD不要構成", "顧客側でSTEP/DXF/PDFへ出力済みのデータを取り込む", "0円", "中\n（材質・質量は落ちる）", "顧客に出力運用を依頼する必要あり"],
            ["④ 変換代行\n（当社の余剰ライセンス活用）", "当社の余剰ライセンス＋当社サーバーでICAD→STEP/中間データ変換を代行", "0円", "最大〜中", "図面データを社外（当社）へ出すことへの顧客同意が必要"],
            ["⑤ 専用ライセンス新規購入", "顧客がナレッジシステム用にICADを1本追加", "約150万円", "最大", "費用が導入判断のボトルネックになりやすい"],
        ],
        col_widths=[1.8, 3.8, 1.5, 1.4, 3],
        font_size=Pt(10),
        header_size=Pt(10.5),
        highlight_rows={1, 4},
        note="当社には余剰ライセンスがあるため、案④は当社が担げる現実的な選択肢です。また対象顧客はICADユーザーであることが多く、案①②が成立する可能性も相応にあると見ています。",
    )

    add_bullet_slide(
        prs,
        "確認が必要な事項（未確認・断定していません）",
        "ライセンス条項の問題であり、当社もまだ確認できていません",
        [
            (0, "富士通／ICAD販売元へ確認する事項", True),
            (1, "設計者向けライセンスを、無人のバッチ処理（サーバー常駐・自動起動）で使うことが許諾範囲か"),
            (1, "ライセンス形態（ノードロック／フローティング等）と、フローティングの場合の借用可否"),
            (1, "「同時使用」の定義。1座席を夜間だけ別PCで使う運用が可能か"),
            (1, "当社の余剰ライセンスを使って他社図面の変換を受託することが許諾範囲か（案④の前提）"),
            (0, "確認結果によって変わること", True),
            (1, "許諾される → 案①②④が使え、顧客の追加ICAD費用ゼロで提案できる"),
            (1, "許諾されない → 案③（ICAD不要構成）を主軸にし、取得情報が減ることを製品仕様として明示する"),
            (0, "ご判断いただきたいこと", True),
            (1, "富士通／販売元への正式な問い合わせを、当社名で行ってよいか"),
        ],
        base_size=Pt(15),
    )

    # CHAPTER 04
    add_chapter_slide(
        prs, "04", "責任分界と\n今後の進め方", "Scope & Next Steps",
        [
            (0, "どこまで自社が作り、どこから創屋が実装するか"),
            (0, "ソース切り出し提供の話"),
            (0, "今後の進め方とご判断事項"),
        ],
    )

    add_table_slide(
        prs,
        "創屋との責任分界",
        "当社は「CADから情報を抜き、タグ・属性にする」まで。本体への組み込みは創屋",
        ["工程", "担当", "現状"],
        [
            ["ICADからの情報抽出（C# + SXNET）", "アルパイン", "実装済み。実ICAD 39件で動作確認済み"],
            ["STEP/DXFからの情報抽出", "アルパイン", "実装済み。外部ライブラリ不要"],
            ["ICAD→STEP/DXF変換", "アルパイン", "実装済み。実機で変換・抽出まで確認済み"],
            ["正規化（属性への整形）", "アルパイン", "実装済み。2D:133項目 / 3D:97項目"],
            ["タグ生成・辞書管理", "アルパイン", "実装済み。辞書はDB化＋画面から編集可能"],
            ["ナレッジシステム本体への登録・更新・削除", "創屋", "創屋側で実装。当社は抽出・正規化・候補生成とAPI契約案を提供"],
            ["2D/3Dビューワーとの連携", "アルパイン→創屋", "当社側でプレビュー生成まで実装。組み込み範囲は要確認"],
            ["検索・RAG・画面", "創屋", "本体側"],
        ],
        col_widths=[4.5, 2, 5],
        font_size=Pt(11.5),
        header_size=Pt(11),
        note="創屋から「完成後はタグ・属性抽出の部分のソースだけ頂ければ、その部分だけ切り出しできると考えている」との要望を受けています。現状の実装はその形で切り出せる構成になっています。",
    )

    add_cards_slide(
        prs,
        "ソース切り出し提供について",
        "創屋の要望「抽出部分のソースだけ欲しい」への回答",
        [
            ("切り出せる構成になっています", [
                "C#側：独立ソリューション（約4,700行）",
                "　ナレッジシステム本体への依存なし",
                "　1図面 = 1回のプロセス呼び出し、入出力はJSON",
                "",
                "Python側：6ファイル（約2,800行）",
                "　Djangoへの結合は実質2点だけ",
                "　（バージョン定数3個 + 辞書テーブル1個）",
                "　→ 純Pythonパッケージとして切り出せる",
            ]),
            ("経営上の論点", [
                "抽出ロジックは当社の技術資産です。",
                "",
                "「どの範囲を、どういう条件で渡すか」は",
                "契約・知財の観点でも決めておく必要があります。",
                "",
                "・ソース提供の対価をどう見るか",
                "・第三者への再利用制限を付けるか",
                "・外販時の当社の取り分をどう設計するか",
                "",
                "→ 技術的には渡せます。条件のご判断をお願いします。",
            ]),
        ],
        cols=2,
        body_size=Pt(12),
    )

    add_table_slide(
        prs,
        "今後の進め方とご判断事項",
        "直近で決めていただきたいことを整理しました",
        ["#", "事項", "当社の考え", "ご判断"],
        [
            ["1", "ICADライセンスの構成方針", "「ICADあり構成」と「ICADなし構成」の両方を成立させる", "この方針で進めてよいか"],
            ["2", "富士通／販売元への確認", "ライセンス条項（無人バッチ・兼用・変換受託の可否）を正式に確認したい", "当社名で問い合わせてよいか"],
            ["3", "抽出ソースの創屋への提供条件", "技術的には切り出して渡せる。契約・知財条件は未定", "提供範囲と条件のご判断"],
            ["4", "社内での試験運用", "定量効果を測るため、自社の実案件で試験運用したい", "対象部署・期間のご判断"],
            ["5", "辞書整備の体制", "客先・装置・案件の語彙登録が精度に直結する。誰が整備するか決めたい", "担当のご判断"],
            ["6", "外販の是非と時期", "自社で運用実績を作ってから判断するのが現実的", "方針のご確認"],
        ],
        col_widths=[0.5, 3, 5, 2.8],
        font_size=Pt(11),
        header_size=Pt(11),
        height=Emu(3200000),
    )

    add_end_slide(prs)
    return finish(prs, OUTDIR / "20260728_ナレッジシステム_タグ属性抽出_役員説明資料_r1.pptx")


if __name__ == "__main__":
    for path in (build_souya_deck(), build_exec_deck()):
        print("created:", path)
